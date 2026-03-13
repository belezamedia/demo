# app_paid_DEV.py
# DocHelp.AI — Upload → Index → Chat (LangChain + Chroma + AWS Bedrock)
# Hybrid retrieval version:
# - semantic search + keyword search across all uploaded chunks
# - safer file parsing
# - better upload status
# - stronger retrieval for all uploaded files
# - more helpful grounded answers without using outside knowledge

import os
import io
import re
import json
import uuid
import time
import shutil
import hmac
import hashlib
import configparser
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Tuple
from urllib.parse import unquote, quote
from html.parser import HTMLParser
from email import policy
from email.parser import BytesParser
from datetime import datetime
import xml.etree.ElementTree as ET
from collections import Counter

import streamlit as st

# Load .env locally only (never override Streamlit Cloud secrets)
try:
    from dotenv import load_dotenv

    if Path(".env").exists():
        load_dotenv(override=False)
except Exception:
    pass

# Core deps
import pandas as pd
import yaml
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain_aws import ChatBedrockConverse, BedrockEmbeddings

# Optional parsers
try:
    from striprtf.striprtf import rtf_to_text as _rtf_to_text
except Exception:
    _rtf_to_text = None

try:
    import extract_msg  # .msg
except Exception:
    extract_msg = None

try:
    import olefile  # .doc / .ppt / .msg fallback
except Exception:
    olefile = None


# ============================
# CONFIG
# ============================
APP_TITLE = "DocHelp.AI"
BRAND_PINK = "#1F4ED8"

DEFAULT_REGION = os.getenv("AWS_DEFAULT_REGION", "us-west-2").strip()

BEDROCK_CHAT_MODEL_ID = unquote(
    os.getenv("BEDROCK_CHAT_MODEL_ID", "us.amazon.nova-pro-v1:0").strip()
)
BEDROCK_EMBED_MODEL_ID = unquote(
    os.getenv("BEDROCK_EMBED_MODEL_ID", "amazon.titan-embed-text-v2:0").strip()
)


def default_app_data_root() -> Path:
    env_root = os.getenv("DOCHELP_DATA_ROOT", "").strip()
    if env_root:
        return Path(env_root).expanduser().resolve()
    return (Path.home() / ".dochelp_ai").resolve()


APP_DATA_ROOT = default_app_data_root()
CHROMA_ROOT = APP_DATA_ROOT / "chroma"

MAX_FILES = int(os.getenv("MAX_FILES", "300"))
MAX_FILE_MB = int(os.getenv("MAX_FILE_MB", "500"))
MAX_TOTAL_UPLOAD_MB = int(os.getenv("MAX_TOTAL_UPLOAD_MB", "3000"))
MAX_TOTAL_CHUNKS = int(os.getenv("MAX_TOTAL_CHUNKS", "100000"))

TOP_K = int(os.getenv("TOP_K", "12"))
SEMANTIC_K = int(os.getenv("SEMANTIC_K", "12"))
KEYWORD_K = int(os.getenv("KEYWORD_K", "12"))
FINAL_K = int(os.getenv("FINAL_K", "14"))

SEMANTIC_K_PER_QUERY = int(os.getenv("SEMANTIC_K_PER_QUERY", str(SEMANTIC_K)))
KEYWORD_K_PER_QUERY = int(os.getenv("KEYWORD_K_PER_QUERY", str(KEYWORD_K)))
ADJACENT_CHUNKS = int(os.getenv("ADJACENT_CHUNKS", "1"))

CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "700"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "120"))
MAX_CONTEXT_CHARS = int(os.getenv("MAX_CONTEXT_CHARS", "25000"))

WORKSPACE_TTL_SECONDS = int(os.getenv("WORKSPACE_TTL_SECONDS", "7200"))  # 2 hours
IDLE_TTL_SECONDS = int(os.getenv("IDLE_TTL_SECONDS", "1800"))  # 30 minutes
JANITOR_MAX_DELETE = int(os.getenv("JANITOR_MAX_DELETE", "50"))
META_FILENAME = "_meta.json"
MANIFEST_FILENAME = "_manifest.json"

TENANT_SIGNING_KEY = os.getenv("TENANT_SIGNING_KEY", "").strip()

APP_USERNAME = os.getenv("APP_USERNAME", "").strip()
APP_PASSWORD = os.getenv("APP_PASSWORD", "").strip()

DEBUG_RETRIEVAL = os.getenv("DEBUG_RETRIEVAL", "false").strip().lower() == "true"

SYSTEM_PROMPT = """You are a professional business assistant.

Rules:
1. Answer using only the provided context from the uploaded files.
2. Do not use outside knowledge.
3. If the answer is partially supported by the uploaded files, provide the best grounded answer and clearly note what is uncertain or missing.
4. If the answer is not supported by the uploaded files, say exactly:
   "I don't know based on the uploaded files."
5. Write in a concise, professional, businesslike tone.
6. Prefer a direct answer first.
7. When useful, mention the file names that support the answer.
8. Never invent policy details, numbers, dates, entities, or conclusions not present in the context.

Return VALID JSON only in this format:
{"answer":"..."}

No markdown. No extra keys. No explanations.
"""

QUERY_REWRITE_PROMPT = """You rewrite user questions into search-friendly variants for document retrieval.

Rules:
- Stay faithful to the original meaning.
- Produce 3 concise alternate queries.
- Include likely key terms, synonyms, and document-oriented phrasing.
- Do not add facts not present in the question.
- If the question mentions a file name, preserve it.
- Return valid JSON only:
{"queries":["q1","q2","q3"]}
"""


# ============================
# UI / STYLES
# ============================
def inject_styles():
    st.markdown(
        f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

:root {{
  --brand-blue: {BRAND_PINK};
}}

html, body, [class*="css"] {{
  font-family: 'Inter', sans-serif;
}}

.gothic-title {{
  font-family: 'Inter', sans-serif;
  font-size: 56px;
  font-weight: 700;
  line-height: 1.0;
  text-align: center;
  margin-top: 0.5rem;
  margin-bottom: 0.35rem;
}}

.gothic-sub {{
  font-family: 'Inter', sans-serif;
  font-size: 20px;
  font-weight: 600;
  text-align: center;
  margin-bottom: 0.75rem;
}}

.center-note {{
  text-align: center;
  font-size: 18px;
  color: rgba(0,0,0,0.72);
  margin-bottom: 1rem;
}}

.small-foot {{
  text-align: center;
  font-size: 13px;
  color: rgba(0,0,0,0.65);
  margin-top: 1rem;
  line-height: 1.65;
}}

.small-foot em {{
  font-style: italic;
}}

.small-muted {{
  text-align: center;
  font-size: 13px;
  color: rgba(0,0,0,0.58);
  margin-top: 0.25rem;
  margin-bottom: 0.5rem;
}}

.small-success {{
  font-size: 13px;
  color: rgba(0,0,0,0.72);
  margin-top: 0.25rem;
  margin-bottom: 0.5rem;
}}

#MainMenu {{visibility: hidden;}}
footer {{visibility: hidden;}}
header {{visibility: hidden;}}

button[kind="primary"] {{
  background: var(--brand-blue) !important;
  border: 1px solid var(--brand-blue) !important;
  color: white !important;
  font-family: 'Inter', sans-serif !important;
}}

div.stButton > button {{
  border-radius: 12px !important;
  padding: 0.65rem 1rem !important;
}}

div[data-testid="stTextInput"] input {{
  border: 2px solid var(--brand-blue) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
}}

div[data-testid="stTextInput"] input:focus {{
  border: 2px solid var(--brand-blue) !important;
  outline: none !important;
  box-shadow: 0 0 0 3px rgba(31, 78, 216, 0.15) !important;
}}

div[data-testid="stChatInput"] textarea {{
  border: 2px solid var(--brand-blue) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
}}

div[data-testid="stChatInput"] textarea:focus {{
  border: 2px solid var(--brand-blue) !important;
  outline: none !important;
  box-shadow: 0 0 0 3px rgba(31, 78, 216, 0.15) !important;
}}

div[data-testid="stChatInput"] button {{
  border-radius: 12px !important;
  border: 1px solid var(--brand-blue) !important;
  background: var(--brand-blue) !important;
}}

div[data-testid="stChatInput"] button svg {{
  fill: white !important;
  stroke: white !important;
}}
</style>
""",
        unsafe_allow_html=True,
    )


# ============================
# LOGIN GATE
# ============================
def login_gate() -> None:
    if st.session_state.get("authed") is True:
        return

    st.markdown(f'<div class="gothic-title">{APP_TITLE}</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="center-note">Private AI for Your Documents</div>',
        unsafe_allow_html=True,
    )

    if not (APP_USERNAME and APP_PASSWORD):
        st.error(
            "Missing APP_USERNAME / APP_PASSWORD. Set them in .env (local) or Streamlit Secrets (cloud)."
        )
        st.stop()

    st.markdown('<div class="gothic-sub">Login</div>', unsafe_allow_html=True)

    u = st.text_input("username", label_visibility="collapsed", placeholder="username")
    p = st.text_input(
        "password",
        type="password",
        label_visibility="collapsed",
        placeholder="password",
    )

    c1, c2 = st.columns([2, 1])
    with c1:
        go = st.button("Enter", type="primary", use_container_width=True)
    with c2:
        if st.button("Reset", use_container_width=True):
            st.session_state.clear()
            st.rerun()

    if go:
        if u == APP_USERNAME and p == APP_PASSWORD:
            st.session_state["authed"] = True
            st.rerun()
        else:
            st.error("Invalid username or password.")

    st.stop()


# ============================
# TENANT ISOLATION
# ============================
def _b(x: str) -> bytes:
    return (x or "").encode("utf-8")


def _is_streamlit_cloud() -> bool:
    return Path("/mount/src").exists()


def sign_token(token: str) -> str:
    if _is_streamlit_cloud() and not TENANT_SIGNING_KEY:
        raise RuntimeError("TENANT_SIGNING_KEY must be set in Streamlit Secrets.")
    key = TENANT_SIGNING_KEY or "insecure-dev-key-change-me"
    return hmac.new(_b(key), _b(token), hashlib.sha256).hexdigest()


def verify_token(token: str, sig: str) -> bool:
    if not token or not sig:
        return False
    expected = sign_token(token)
    return hmac.compare_digest(expected, sig)


def workspace_id_from_token(token: str) -> str:
    return hashlib.sha256(_b(token)).hexdigest()


def set_query_params(token: str, sig: str) -> None:
    try:
        st.query_params["t"] = token
        st.query_params["sig"] = sig
    except Exception:
        try:
            st.experimental_set_query_params(t=token, sig=sig)
        except Exception:
            pass


def get_query_params() -> Dict[str, str]:
    try:
        qp = st.query_params
        return {"t": qp.get("t", ""), "sig": qp.get("sig", "")}
    except Exception:
        qp = st.experimental_get_query_params()
        return {
            "t": (qp.get("t", [""]) or [""])[0],
            "sig": (qp.get("sig", [""]) or [""])[0],
        }


def ensure_tenant_context() -> str:
    qp = get_query_params()
    t = (qp.get("t") or "").strip()
    sig = (qp.get("sig") or "").strip()

    valid = False
    try:
        valid = verify_token(t, sig)
    except Exception:
        valid = False

    if not valid:
        t = uuid.uuid4().hex
        sig = sign_token(t)
        set_query_params(t, sig)

    ws = workspace_id_from_token(t)
    st.session_state["workspace_id"] = ws
    return ws


# ============================
# AVATAR
# ============================
def assistant_avatar_data_uri(color_hex: str) -> str:
    svg = f"""
    <svg xmlns="http://www.w3.org/2000/svg" width="64" height="64" viewBox="0 0 64 64">
      <rect x="14" y="18" width="36" height="32" rx="10" fill="{color_hex}"/>
      <rect x="26" y="10" width="12" height="10" rx="4" fill="{color_hex}"/>
      <circle cx="26" cy="34" r="4" fill="white"/>
      <circle cx="38" cy="34" r="4" fill="white"/>
      <rect x="24" y="42" width="16" height="4" rx="2" fill="white" opacity="0.95"/>
      <circle cx="32" cy="8" r="3" fill="{color_hex}"/>
      <rect x="31" y="8" width="2" height="6" fill="{color_hex}"/>
    </svg>
    """.strip()
    return "data:image/svg+xml;utf8," + quote(svg)


ASSISTANT_AVATAR = assistant_avatar_data_uri(BRAND_PINK)


# ============================
# CACHED RESOURCES
# ============================
@st.cache_resource(show_spinner=False)
def get_splitter() -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", "; ", ", ", " ", ""],
    )


@st.cache_resource(show_spinner=False)
def get_embeddings() -> BedrockEmbeddings:
    return BedrockEmbeddings(
        region_name=DEFAULT_REGION,
        model_id=BEDROCK_EMBED_MODEL_ID,
    )


@st.cache_resource(show_spinner=False)
def get_llm() -> ChatBedrockConverse:
    return ChatBedrockConverse(
        model=BEDROCK_CHAT_MODEL_ID,
        region_name=DEFAULT_REGION,
        temperature=0.0,
        max_tokens=700,
    )


# ============================
# BEDROCK STREAMING
# ============================
def _extract_text(obj: Any) -> str:
    if isinstance(obj, str):
        return obj
    if hasattr(obj, "content"):
        return _extract_text(obj.content)
    if isinstance(obj, list):
        parts: List[str] = []
        for x in obj:
            if isinstance(x, dict):
                parts.append(str(x.get("text", "")))
            elif hasattr(x, "get"):
                parts.append(str(x.get("text", "")))
            else:
                parts.append(str(x))
        return "".join(parts)
    return ""


def generate_streaming(prompt: str) -> str:
    llm = get_llm()
    parts: List[str] = []
    try:
        for chunk in llm.stream(prompt):
            txt = _extract_text(chunk)
            if txt:
                parts.append(txt)
        return "".join(parts)
    except Exception:
        resp = llm.invoke(prompt)
        return _extract_text(resp)


def safe_json_load(s: str) -> Dict[str, Any]:
    try:
        obj = json.loads(s)
        if isinstance(obj, dict):
            return obj
    except Exception:
        pass

    m = re.search(r"\{.*\}", s, flags=re.DOTALL)
    if m:
        try:
            obj = json.loads(m.group(0))
            if isinstance(obj, dict):
                return obj
        except Exception:
            pass

    return {"answer": s.strip()}


def rewrite_queries(question: str) -> List[str]:
    prompt = f"""{QUERY_REWRITE_PROMPT}

Original question:
{question}
"""
    try:
        raw = generate_streaming(prompt)
        obj = safe_json_load(raw)
        queries = obj.get("queries", [])
        if isinstance(queries, list):
            clean = []
            for q in queries:
                q = str(q or "").strip()
                if q and q.lower() != question.lower():
                    clean.append(q)
            return dedupe_preserve_order([question] + clean[:3])
    except Exception:
        pass
    return [question]


# ============================
# FILE PARSING
# ============================
class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self._chunks: List[str] = []

    def handle_data(self, data: str) -> None:
        if data and data.strip():
            self._chunks.append(data.strip())

    def get_text(self) -> str:
        return " ".join(self._chunks).strip()


def _clean_whitespace(s: str) -> str:
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\r\n?", "\n", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _decode_best_effort(file_bytes: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except Exception:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def _strip_xml_tags(xml_text: str) -> str:
    try:
        root = ET.fromstring(xml_text)
        parts: List[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        return _clean_whitespace("\n".join(parts))
    except Exception:
        return _clean_whitespace(re.sub(r"<[^>]+>", " ", xml_text))


def _read_zip_member_text(file_bytes: bytes, member_names: List[str]) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            parts: List[str] = []
            for name in member_names:
                if name in zf.namelist():
                    xml_text = zf.read(name).decode("utf-8", errors="ignore")
                    cleaned = _strip_xml_tags(xml_text)
                    if cleaned:
                        parts.append(cleaned)
            return _clean_whitespace("\n\n".join(parts))
    except Exception:
        return ""
    return ""


def read_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    texts: List[str] = []
    for i, p in enumerate(reader.pages):
        try:
            txt = p.extract_text() or ""
        except Exception:
            txt = ""
        if txt.strip():
            texts.append(f"[Page {i + 1}]\n{txt}")
    return _clean_whitespace("\n\n".join(texts))


def read_csv(file_bytes: bytes, sep: str = ",") -> str:
    last_err = None
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            df = pd.read_csv(io.BytesIO(file_bytes), sep=sep, encoding=enc)
            return df.to_csv(index=False)
        except Exception as e:
            last_err = e
    raise RuntimeError(f"Could not read CSV/TSV/PSV file: {last_err}")


def read_excel(file_bytes: bytes, filename: str = "") -> str:
    name = (filename or "").lower()
    bio = io.BytesIO(file_bytes)
    try:
        if name.endswith(".xlsx"):
            sheets = pd.read_excel(bio, engine="openpyxl", sheet_name=None)
        elif name.endswith(".xls"):
            sheets = pd.read_excel(bio, engine="xlrd", sheet_name=None)
        else:
            sheets = pd.read_excel(bio, sheet_name=None)
    except ImportError as e:
        raise RuntimeError(
            f"Excel dependency missing for '{filename}'. Make sure openpyxl is installed for .xlsx and xlrd for .xls."
        ) from e
    except Exception as e:
        raise RuntimeError(f"Could not read Excel file '{filename}': {e}") from e

    parts: List[str] = []
    if isinstance(sheets, dict):
        for sheet_name, df in sheets.items():
            parts.append(f"=== Sheet: {sheet_name} ===")
            try:
                parts.append(df.fillna("").to_csv(index=False))
            except Exception:
                parts.append(str(df))
            parts.append("")
    else:
        parts.append(str(sheets))
    return _clean_whitespace("\n".join(parts))


def read_xlsb(file_bytes: bytes, filename: str = "") -> str:
    try:
        sheets = pd.read_excel(io.BytesIO(file_bytes), engine="pyxlsb", sheet_name=None)
        parts: List[str] = []
        if isinstance(sheets, dict):
            for sheet_name, df in sheets.items():
                parts.append(f"=== Sheet: {sheet_name} ===")
                parts.append(df.fillna("").to_csv(index=False))
                parts.append("")
        return _clean_whitespace("\n".join(parts))
    except ImportError as e:
        raise RuntimeError(
            f"Excel dependency missing for '{filename}'. Make sure pyxlsb is installed for .xlsb files."
        ) from e
    except Exception as e:
        raise RuntimeError(f"Could not read XLSB file '{filename}': {e}") from e


def read_parquet(file_bytes: bytes, filename: str = "") -> str:
    try:
        df = pd.read_parquet(io.BytesIO(file_bytes))
        return df.fillna("").to_csv(index=False)
    except Exception as e:
        raise RuntimeError(f"Could not read Parquet file '{filename}': {e}") from e


def read_json(file_bytes: bytes) -> str:
    try:
        obj = json.loads(_decode_best_effort(file_bytes))
        return json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception:
        return _decode_best_effort(file_bytes).strip()


def read_jsonl(file_bytes: bytes) -> str:
    text = _decode_best_effort(file_bytes)
    rows: List[str] = []
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            obj = json.loads(line)
            rows.append(json.dumps(obj, ensure_ascii=False))
        except Exception:
            rows.append(line)
    return _clean_whitespace("\n".join(rows))


def read_yaml(file_bytes: bytes) -> str:
    text = _decode_best_effort(file_bytes)
    try:
        obj = yaml.safe_load(text)
        return yaml.safe_dump(obj, sort_keys=False, allow_unicode=True)
    except Exception:
        return text.strip()


def read_ini(file_bytes: bytes) -> str:
    raw = _decode_best_effort(file_bytes)
    cp = configparser.ConfigParser()
    try:
        cp.read_string(raw)
        out_lines: List[str] = []
        for section in cp.sections():
            out_lines.append(f"[{section}]")
            for k, v in cp.items(section):
                out_lines.append(f"{k} = {v}")
            out_lines.append("")
        return _clean_whitespace("\n".join(out_lines))
    except Exception:
        return raw.strip()


def read_xml(file_bytes: bytes) -> str:
    try:
        root = ET.fromstring(file_bytes)
        parts: List[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        return _clean_whitespace("\n".join(parts))
    except Exception:
        xml = _decode_best_effort(file_bytes)
        return _clean_whitespace(re.sub(r"<[^>]+>", " ", xml))


def read_html(file_bytes: bytes) -> str:
    html = _decode_best_effort(file_bytes)
    parser = _HTMLTextExtractor()
    try:
        parser.feed(html)
        return _clean_whitespace(parser.get_text())
    except Exception:
        return _clean_whitespace(re.sub(r"<[^>]+>", " ", html))


def read_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return _clean_whitespace("\n".join(parts))


def read_pptx(file_bytes: bytes) -> str:
    pres = Presentation(io.BytesIO(file_bytes))
    parts: List[str] = []
    for slide_idx, slide in enumerate(pres.slides, start=1):
        slide_parts: List[str] = []
        for shape in slide.shapes:
            txt = getattr(shape, "text", None)
            if txt and str(txt).strip():
                slide_parts.append(str(txt).strip())
        if slide_parts:
            parts.append(f"[Slide {slide_idx}]")
            parts.extend(slide_parts)
            parts.append("")
    return _clean_whitespace("\n".join(parts))


def read_odt(file_bytes: bytes) -> str:
    return _read_zip_member_text(file_bytes, ["content.xml", "styles.xml", "meta.xml"])


def read_ods(file_bytes: bytes) -> str:
    return _read_zip_member_text(file_bytes, ["content.xml", "styles.xml", "meta.xml"])


def read_odp(file_bytes: bytes) -> str:
    return _read_zip_member_text(file_bytes, ["content.xml", "styles.xml", "meta.xml"])


def _extract_printable_utf16le(blob: bytes) -> List[str]:
    text = blob.decode("utf-16le", errors="ignore")
    candidates = re.findall(r"[\x20-\x7E][\x20-\x7E\n\r\t]{3,}", text)
    return [c.strip() for c in candidates if c.strip()]


def _extract_printable_utf8(blob: bytes) -> List[str]:
    text = blob.decode("utf-8", errors="ignore")
    candidates = re.findall(r"[A-Za-z0-9][^\x00]{3,}", text)
    return [c.strip() for c in candidates if c.strip()]


def read_doc_legacy(file_bytes: bytes) -> str:
    if olefile is None:
        return ""
    try:
        with olefile.OleFileIO(io.BytesIO(file_bytes)) as ole:
            pieces: List[str] = []
            for stream_name in ole.listdir():
                joined = "/".join(stream_name)
                if joined in ("WordDocument", "1Table", "0Table"):
                    try:
                        blob = ole.openstream(stream_name).read()
                        pieces.extend(_extract_printable_utf16le(blob))
                        pieces.extend(_extract_printable_utf8(blob))
                    except Exception:
                        continue
            return _clean_whitespace("\n".join(dict.fromkeys(pieces)))
    except Exception:
        return ""


def read_ppt_legacy(file_bytes: bytes) -> str:
    if olefile is None:
        return ""
    try:
        with olefile.OleFileIO(io.BytesIO(file_bytes)) as ole:
            pieces: List[str] = []
            for stream_name in ole.listdir():
                joined = "/".join(stream_name)
                if joined == "PowerPoint Document":
                    try:
                        blob = ole.openstream(stream_name).read()
                        pieces.extend(_extract_printable_utf16le(blob))
                        pieces.extend(_extract_printable_utf8(blob))
                    except Exception:
                        continue
            return _clean_whitespace("\n".join(dict.fromkeys(pieces)))
    except Exception:
        return ""


def read_msg(file_bytes: bytes) -> str:
    if extract_msg is not None:
        tmp_path = None
        try:
            tmp_dir = Path("/tmp")
            tmp_dir.mkdir(parents=True, exist_ok=True)
            tmp_path = tmp_dir / f"msg_{uuid.uuid4().hex}.msg"
            tmp_path.write_bytes(file_bytes)

            msg = extract_msg.Message(str(tmp_path))
            parts: List[str] = []
            if getattr(msg, "subject", None):
                parts.append(f"Subject: {msg.subject}")
            if getattr(msg, "sender", None):
                parts.append(f"From: {msg.sender}")
            if getattr(msg, "to", None):
                parts.append(f"To: {msg.to}")
            if getattr(msg, "date", None):
                parts.append(f"Date: {msg.date}")
            body = getattr(msg, "body", "") or ""
            if body.strip():
                parts.append("")
                parts.append(body.strip())
            return _clean_whitespace("\n".join(parts))
        except Exception:
            pass
        finally:
            try:
                if tmp_path and tmp_path.exists():
                    tmp_path.unlink()
            except Exception:
                pass

    if olefile is None:
        return ""

    try:
        with olefile.OleFileIO(io.BytesIO(file_bytes)) as ole:
            parts: List[str] = []
            for stream_name in ole.listdir():
                joined = "/".join(stream_name)
                if "__substg1.0_" in joined:
                    try:
                        blob = ole.openstream(stream_name).read()
                        parts.extend(_extract_printable_utf16le(blob))
                    except Exception:
                        continue
            return _clean_whitespace("\n".join(dict.fromkeys(parts)))
    except Exception:
        return ""


def read_ipynb(file_bytes: bytes) -> str:
    try:
        nb = json.loads(_decode_best_effort(file_bytes))
        cells = nb.get("cells", []) if isinstance(nb, dict) else []
        parts: List[str] = []
        for c in cells:
            src = c.get("source", [])
            if isinstance(src, list):
                s = "".join(src)
            else:
                s = str(src or "")
            if s.strip():
                parts.append(s.strip())
        return _clean_whitespace("\n\n".join(parts))
    except Exception:
        return _decode_best_effort(file_bytes).strip()


def read_eml(file_bytes: bytes) -> str:
    try:
        msg = BytesParser(policy=policy.default).parsebytes(file_bytes)
        parts: List[str] = []
        subj = msg.get("subject", "")
        frm = msg.get("from", "")
        to = msg.get("to", "")
        if subj:
            parts.append(f"Subject: {subj}")
        if frm:
            parts.append(f"From: {frm}")
        if to:
            parts.append(f"To: {to}")
        parts.append("")

        body_texts: List[str] = []
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                disp = part.get_content_disposition()
                if disp == "attachment":
                    continue
                if ctype in ("text/plain", "text/html"):
                    try:
                        content = part.get_content()
                    except Exception:
                        content = part.get_payload(decode=True)
                        content = (content or b"").decode("utf-8", errors="ignore")
                    if ctype == "text/html":
                        content = read_html(str(content).encode("utf-8", errors="ignore"))
                    if content and str(content).strip():
                        body_texts.append(str(content).strip())
        else:
            ctype = msg.get_content_type()
            content = msg.get_content()
            if ctype == "text/html":
                content = read_html(str(content).encode("utf-8", errors="ignore"))
            if content and str(content).strip():
                body_texts.append(str(content).strip())

        if body_texts:
            parts.append("\n\n".join(body_texts))

        return _clean_whitespace("\n".join(parts))
    except Exception:
        return _decode_best_effort(file_bytes).strip()


def read_rtf(file_bytes: bytes) -> str:
    raw = _decode_best_effort(file_bytes)
    if _rtf_to_text is not None:
        try:
            return _clean_whitespace(_rtf_to_text(raw))
        except Exception:
            pass
    fallback = re.sub(r"{\\.*?}|\\[a-zA-Z]+\d* ?", " ", raw)
    return _clean_whitespace(fallback)


def read_txt(file_bytes: bytes) -> str:
    return _clean_whitespace(_decode_best_effort(file_bytes))


CODE_EXTS = (
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cpp", ".h", ".hpp", ".go",
    ".rs", ".rb", ".php", ".cs", ".swift", ".kt", ".scala", ".sh", ".bash", ".zsh",
    ".ps1", ".sql", ".toml", ".env", ".dockerfile"
)
TEXTY_EXTS = (
    ".txt", ".md", ".log", ".cfg", ".conf", ".ini", ".rst", ".tex"
)


def file_to_text(uploaded_file) -> str:
    name = (uploaded_file.name or "").lower()
    b = uploaded_file.getvalue()

    if name.endswith(".pdf"):
        return read_pdf(b)

    if name.endswith(".docx"):
        return read_docx(b)
    if name.endswith(".doc"):
        return read_doc_legacy(b)

    if name.endswith(".pptx"):
        return read_pptx(b)
    if name.endswith(".ppt"):
        return read_ppt_legacy(b)

    if name.endswith(".odt"):
        return read_odt(b)
    if name.endswith(".ods"):
        return read_ods(b)
    if name.endswith(".odp"):
        return read_odp(b)

    if name.endswith(".csv"):
        return read_csv(b, sep=",")
    if name.endswith(".tsv"):
        return read_csv(b, sep="\t")
    if name.endswith(".psv"):
        return read_csv(b, sep="|")
    if name.endswith(".xlsb"):
        return read_xlsb(b, name)
    if name.endswith((".xls", ".xlsx")):
        return read_excel(b, name)
    if name.endswith(".parquet"):
        return read_parquet(b, name)

    if name.endswith((".html", ".htm")):
        return read_html(b)
    if name.endswith(".xml"):
        return read_xml(b)
    if name.endswith((".yaml", ".yml")):
        return read_yaml(b)
    if name.endswith((".ini", ".cfg", ".conf")):
        return read_ini(b)

    if name.endswith(".ipynb"):
        return read_ipynb(b)
    if name.endswith(".eml"):
        return read_eml(b)
    if name.endswith(".msg"):
        return read_msg(b)
    if name.endswith(".rtf"):
        return read_rtf(b)

    if name.endswith(".json"):
        return read_json(b)
    if name.endswith((".jsonl", ".ndjson")):
        return read_jsonl(b)

    if name.endswith(CODE_EXTS) or name.endswith(TEXTY_EXTS):
        return read_txt(b)

    return ""


def validate_uploads(files) -> Tuple[bool, str]:
    if not files:
        return False, "Please upload at least one file."

    if len(files) > MAX_FILES:
        return False, f"Too many files. Max is {MAX_FILES}."

    total_bytes = 0
    for f in files:
        size_bytes = int(getattr(f, "size", 0) or 0)
        total_bytes += size_bytes

        size_mb = size_bytes / (1024 * 1024)
        if size_mb > MAX_FILE_MB:
            return False, f"File '{f.name}' is too large ({size_mb:.1f} MB). Max is {MAX_FILE_MB} MB."

    total_mb = total_bytes / (1024 * 1024)
    if total_mb > MAX_TOTAL_UPLOAD_MB:
        return False, f"Total upload too large ({total_mb:.1f} MB). Max total is {MAX_TOTAL_UPLOAD_MB} MB."

    return True, ""


# ============================
# WORKSPACE META + JANITOR
# ============================
def _now() -> int:
    return int(time.time())


def utc_stamp() -> str:
    return datetime.utcnow().strftime("%Y%m%d%H%M%S%f")


def ensure_dir(path: Path) -> Path:
    path.mkdir(parents=True, exist_ok=True)
    return path


def assert_writable_dir(path: Path) -> None:
    ensure_dir(path)
    test_file = path / f".write_test_{uuid.uuid4().hex}.tmp"
    try:
        test_file.write_text("ok", encoding="utf-8")
        _ = test_file.read_text(encoding="utf-8")
    finally:
        try:
            if test_file.exists():
                test_file.unlink()
        except Exception:
            pass


def safe_rmtree(path: Path) -> None:
    try:
        if path.exists():
            shutil.rmtree(path, ignore_errors=True)
    except Exception:
        pass


def read_workspace_meta(d: Path) -> Dict[str, Any]:
    p = d / META_FILENAME
    try:
        return json.loads(p.read_text(encoding="utf-8"))
    except Exception:
        return {}


def write_workspace_meta(d: Path, created_at: int, last_activity: int) -> None:
    ensure_dir(d)
    p = d / META_FILENAME
    meta = {"created_at": int(created_at), "last_activity": int(last_activity)}
    try:
        p.write_text(json.dumps(meta), encoding="utf-8")
    except Exception:
        pass


def get_workspace_db_version() -> str:
    version = st.session_state.get("db_version")
    if not version:
        version = utc_stamp()
        st.session_state["db_version"] = version
    return version


def bump_workspace_db_version() -> str:
    version = utc_stamp()
    st.session_state["db_version"] = version
    return version


def workspace_base_dir(workspace_id: str) -> Path:
    return ensure_dir(CHROMA_ROOT / workspace_id)


def workspace_meta_dir(workspace_id: str) -> Path:
    d = workspace_base_dir(workspace_id) / "meta"
    ensure_dir(d)
    assert_writable_dir(d)
    return d


def workspace_db_dir(workspace_id: str) -> Path:
    d = workspace_base_dir(workspace_id) / f"db_{get_workspace_db_version()}"
    ensure_dir(d)
    assert_writable_dir(d)
    return d


def janitor_sweep() -> int:
    if not CHROMA_ROOT.exists():
        return 0

    now = _now()
    deleted = 0

    for ws_dir in CHROMA_ROOT.iterdir():
        if not ws_dir.is_dir():
            continue

        meta_dir = ws_dir / "meta"
        meta = read_workspace_meta(meta_dir) if meta_dir.exists() else {}
        created_at = int(meta.get("created_at") or 0)
        last_activity = int(meta.get("last_activity") or 0)
        mtime = int(ws_dir.stat().st_mtime)

        if created_at <= 0:
            created_at = mtime
        if last_activity <= 0:
            last_activity = mtime

        if (now - created_at) > WORKSPACE_TTL_SECONDS or (now - last_activity) > IDLE_TTL_SECONDS:
            safe_rmtree(ws_dir)
            deleted += 1

        if deleted >= JANITOR_MAX_DELETE:
            break

    return deleted


def touch_activity(workspace_id: str) -> None:
    now = _now()

    if "created_at" not in st.session_state:
        st.session_state.created_at = now

    st.session_state.last_activity = now
    write_workspace_meta(
        workspace_meta_dir(workspace_id),
        st.session_state.created_at,
        st.session_state.last_activity,
    )


# ============================
# VECTOR STORE + LOCAL CHUNK CATALOG
# ============================
def workspace_dir(workspace_id: str) -> Path:
    d = workspace_meta_dir(workspace_id)
    created_at = int(st.session_state.get("created_at") or _now())
    last_activity = int(st.session_state.get("last_activity") or _now())
    write_workspace_meta(d, created_at, last_activity)
    return d


def manifest_path(workspace_id: str) -> Path:
    return workspace_meta_dir(workspace_id) / MANIFEST_FILENAME


def save_manifest(workspace_id: str, manifest: Dict[str, Any]) -> None:
    p = manifest_path(workspace_id)
    ensure_dir(p.parent)
    assert_writable_dir(p.parent)
    p.write_text(json.dumps(manifest, ensure_ascii=False), encoding="utf-8")


def load_manifest(workspace_id: str) -> Dict[str, Any]:
    p = manifest_path(workspace_id)
    try:
        data = json.loads(p.read_text(encoding="utf-8"))
        if isinstance(data, dict):
            data.setdefault("files", [])
            data.setdefault("chunks", [])
            return data
    except Exception:
        pass
    return {"files": [], "chunks": []}


def clear_workspace_storage(workspace_id: str) -> None:
    bump_workspace_db_version()
    write_workspace_meta(
        workspace_meta_dir(workspace_id),
        int(st.session_state.get("created_at") or _now()),
        int(st.session_state.get("last_activity") or _now()),
    )
    save_manifest(workspace_id, {"files": [], "chunks": []})


def get_vectordb(workspace_id: str) -> Chroma:
    db_dir = workspace_db_dir(workspace_id)
    collection = f"demo_docs_{workspace_id[:12]}_{get_workspace_db_version()[:12]}"
    return Chroma(
        collection_name=collection,
        persist_directory=str(db_dir),
        embedding_function=get_embeddings(),
    )


def add_documents_with_retry(workspace_id: str, docs: List[Document], ids: List[str]) -> None:
    if not docs:
        return

    last_error = None
    for attempt in range(2):
        try:
            vectordb = get_vectordb(workspace_id)
            vectordb.add_documents(docs, ids=ids)
            return
        except Exception as e:
            last_error = e
            if attempt == 0:
                bump_workspace_db_version()
                continue

    raise RuntimeError(
        f"Chroma write failed after retry. Current db dir: {workspace_db_dir(workspace_id)}. Error: {last_error}"
    ) from last_error


def reset_workspace() -> None:
    t = uuid.uuid4().hex
    sig = sign_token(t)
    set_query_params(t, sig)

    st.session_state.indexed = False
    st.session_state.last_index_count = 0
    st.session_state.messages = []
    st.session_state.failed_files = []
    st.session_state.processed_files = []
    st.session_state.created_at = _now()
    st.session_state.last_activity = _now()
    st.session_state.last_debug_docs = []
    st.session_state.db_version = utc_stamp()


# ============================
# INDEXING
# ============================
STOPWORDS = {
    "the", "a", "an", "and", "or", "but", "if", "then", "of", "to", "in", "on", "for",
    "with", "by", "from", "at", "as", "is", "are", "was", "were", "be", "been", "being",
    "this", "that", "these", "those", "it", "its", "their", "there", "here", "what",
    "which", "who", "whom", "how", "when", "where", "why", "can", "could", "should",
    "would", "do", "does", "did", "about", "into", "than", "them", "they", "you", "your",
    "me", "my", "we", "our", "please", "show", "tell", "give"
}


def _normalize_text_for_search(s: str) -> str:
    s = (s or "").lower()
    s = re.sub(r"[^a-z0-9\s._/-]", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def tokenize(s: str) -> List[str]:
    s = _normalize_text_for_search(s)
    return [t for t in s.split() if t and t not in STOPWORDS and len(t) > 1]


def dedupe_preserve_order(items: List[str]) -> List[str]:
    out = []
    seen = set()
    for x in items:
        if x not in seen:
            seen.add(x)
            out.append(x)
    return out


def file_hash_bytes(file_bytes: bytes) -> str:
    return hashlib.sha256(file_bytes).hexdigest()


def build_chunk_id(file_hash: str, chunk_idx: int) -> str:
    return f"{file_hash}::{chunk_idx}"


def get_existing_file_hashes(manifest: Dict[str, Any]) -> set:
    return {str(f.get("file_hash")) for f in manifest.get("files", []) if f.get("file_hash")}


def index_files(workspace_id: str, files) -> int:
    clear_workspace_storage(workspace_id)

    assert_writable_dir(APP_DATA_ROOT)
    assert_writable_dir(CHROMA_ROOT)
    assert_writable_dir(workspace_base_dir(workspace_id))
    assert_writable_dir(workspace_meta_dir(workspace_id))
    assert_writable_dir(workspace_db_dir(workspace_id))

    splitter = get_splitter()
    manifest = load_manifest(workspace_id)
    docs: List[Document] = []
    ids: List[str] = []
    failed_files: List[str] = []
    processed_files: List[str] = []
    existing_hashes = get_existing_file_hashes(manifest)

    total_chunks = 0

    for f in files:
        try:
            raw_bytes = f.getvalue()
            f_hash = file_hash_bytes(raw_bytes)
        except Exception as e:
            failed_files.append(f"{f.name}: could not read uploaded bytes: {e}")
            continue

        if f_hash in existing_hashes:
            continue

        try:
            text = file_to_text(f)
        except Exception as e:
            failed_files.append(f"{f.name}: {type(e).__name__}: {e}")
            continue

        if not text.strip():
            failed_files.append(f"{f.name}: no extractable text found")
            continue

        try:
            chunks = splitter.split_text(text)
        except Exception as e:
            failed_files.append(f"{f.name}: chunking failed: {e}")
            continue

        clean_chunks = [(chunk or "").strip() for chunk in chunks if (chunk or "").strip()]
        if not clean_chunks:
            failed_files.append(f"{f.name}: no usable chunks after parsing")
            continue

        if total_chunks + len(clean_chunks) > MAX_TOTAL_CHUNKS:
            failed_files.append("Chunk limit reached; some later file content may not have been indexed.")
            break

        manifest["files"].append(
            {
                "file_name": f.name,
                "file_hash": f_hash,
                "size_bytes": int(getattr(f, "size", 0) or 0),
                "chunk_count": len(clean_chunks),
                "indexed_at": _now(),
            }
        )

        for idx, chunk in enumerate(clean_chunks):
            chunk_id = build_chunk_id(f_hash, idx)
            manifest["chunks"].append(
                {
                    "id": chunk_id,
                    "source": f.name,
                    "file_hash": f_hash,
                    "chunk": idx,
                    "text": chunk,
                    "norm_text": _normalize_text_for_search(chunk),
                }
            )
            docs.append(
                Document(
                    page_content=chunk,
                    metadata={"id": chunk_id, "source": f.name, "file_hash": f_hash, "chunk": idx},
                )
            )
            ids.append(chunk_id)
            total_chunks += 1

        processed_files.append(f.name)
        existing_hashes.add(f_hash)

    if docs:
        add_documents_with_retry(workspace_id, docs, ids)

    save_manifest(workspace_id, manifest)

    st.session_state.indexed = len(manifest.get("chunks", [])) > 0
    st.session_state.last_index_count = len(manifest.get("chunks", []))
    st.session_state.failed_files = dedupe_preserve_order(failed_files)
    st.session_state.processed_files = dedupe_preserve_order(processed_files)
    touch_activity(workspace_id)
    return len(manifest.get("chunks", []))


# ============================
# RETRIEVAL
# ============================
def score_keyword_match(question: str, chunk_text: str, filename: str = "") -> float:
    q_tokens = tokenize(question)
    if not q_tokens:
        return 0.0

    text_norm = _normalize_text_for_search(chunk_text)
    file_norm = _normalize_text_for_search(filename)
    text_tokens = text_norm.split()

    if not text_tokens and not file_norm:
        return 0.0

    q_counter = Counter(q_tokens)
    t_counter = Counter(text_tokens)

    overlap = 0.0
    for tok, q_count in q_counter.items():
        overlap += min(q_count, t_counter.get(tok, 0))

    phrase_bonus = 0.0
    q_norm = _normalize_text_for_search(question)
    if q_norm and q_norm in text_norm:
        phrase_bonus += 5.0

    exact_token_bonus = 0.0
    filename_bonus = 0.0
    for tok in set(q_tokens):
        if tok in text_norm:
            exact_token_bonus += 0.25
        if tok in file_norm:
            filename_bonus += 0.80

    density = overlap / max(1.0, len(q_tokens))
    return float(density * 10.0 + phrase_bonus + exact_token_bonus + filename_bonus)


def keyword_search(workspace_id: str, question: str, k: int = KEYWORD_K) -> List[Document]:
    manifest = load_manifest(workspace_id)
    records = manifest.get("chunks", [])
    if not records:
        return []

    scored: List[Tuple[float, Dict[str, Any]]] = []
    for r in records:
        score = score_keyword_match(question, r.get("text", ""), r.get("source", ""))
        if score > 0:
            scored.append((score, r))

    scored.sort(key=lambda x: x[0], reverse=True)

    docs: List[Document] = []
    for score, r in scored[:k]:
        docs.append(
            Document(
                page_content=r.get("text", ""),
                metadata={
                    "id": r.get("id"),
                    "source": r.get("source", "unknown"),
                    "file_hash": r.get("file_hash"),
                    "chunk": r.get("chunk", 0),
                    "retrieval": "keyword",
                    "keyword_score": round(score, 4),
                },
            )
        )
    return docs


def semantic_search(workspace_id: str, question: str, k: int = SEMANTIC_K) -> List[Document]:
    vectordb = get_vectordb(workspace_id)
    try:
        docs = vectordb.max_marginal_relevance_search(question, k=k, fetch_k=max(k * 3, 24))
    except Exception:
        try:
            docs = vectordb.similarity_search(question, k=k)
        except Exception:
            docs = []

    out: List[Document] = []
    for rank, d in enumerate(docs):
        md = dict(d.metadata or {})
        md["retrieval"] = "semantic"
        md["semantic_score"] = max(0.0, float(k - rank))
        out.append(Document(page_content=d.page_content, metadata=md))
    return out


def merge_results(semantic_docs: List[Document], keyword_docs: List[Document], question: str, final_k: int = FINAL_K) -> List[Document]:
    by_id: Dict[str, Document] = {}
    combined_scores: Dict[str, float] = {}
    retrievals: Dict[str, set] = {}

    for d in semantic_docs + keyword_docs:
        md = dict(d.metadata or {})
        doc_id = str(md.get("id") or f"{md.get('source')}::{md.get('chunk')}")
        score = float(md.get("semantic_score", 0.0) or md.get("keyword_score", 0.0) or 0.0)

        if doc_id not in by_id:
            by_id[doc_id] = d
            combined_scores[doc_id] = score
            retrievals[doc_id] = {md.get("retrieval", "")}
        else:
            combined_scores[doc_id] += score
            retrievals[doc_id].add(md.get("retrieval", ""))

    merged: List[Document] = []
    for doc_id, d in by_id.items():
        md = dict(d.metadata or {})
        md["retrievals"] = sorted(list(retrievals.get(doc_id, set())))
        md["combined_score"] = combined_scores.get(doc_id, 0.0)
        md["final_score"] = (
            float(md["combined_score"])
            + score_keyword_match(question, d.page_content, md.get("source", ""))
            + (1.25 if len(md["retrievals"]) > 1 else 0.0)
        )
        merged.append(Document(page_content=d.page_content, metadata=md))

    merged.sort(key=lambda d: float(d.metadata.get("final_score", 0.0)), reverse=True)

    final_docs: List[Document] = []
    seen = set()
    for d in merged:
        key = (d.metadata.get("source"), d.metadata.get("chunk"))
        if key in seen:
            continue
        seen.add(key)
        final_docs.append(d)
        if len(final_docs) >= final_k:
            break

    return final_docs


def expand_adjacent_docs(workspace_id: str, docs: List[Document]) -> List[Document]:
    manifest = load_manifest(workspace_id)
    records = manifest.get("chunks", [])
    if not records or not docs or ADJACENT_CHUNKS <= 0:
        return docs

    chunk_lookup = {
        (str(r.get("file_hash")), int(r.get("chunk", 0))): r for r in records
    }

    expanded: List[Document] = []
    seen = set()

    for d in docs:
        md = dict(d.metadata or {})
        key = (md.get("source"), md.get("chunk"))
        if key not in seen:
            expanded.append(d)
            seen.add(key)

        file_hash = str(md.get("file_hash"))
        chunk_idx = int(md.get("chunk", 0))
        for offset in range(-ADJACENT_CHUNKS, ADJACENT_CHUNKS + 1):
            if offset == 0:
                continue
            neighbor = chunk_lookup.get((file_hash, chunk_idx + offset))
            if not neighbor:
                continue
            nkey = (neighbor.get("source"), neighbor.get("chunk"))
            if nkey in seen:
                continue
            seen.add(nkey)
            expanded.append(
                Document(
                    page_content=neighbor.get("text", ""),
                    metadata={
                        "id": neighbor.get("id"),
                        "source": neighbor.get("source", "unknown"),
                        "file_hash": neighbor.get("file_hash"),
                        "chunk": neighbor.get("chunk", 0),
                        "retrieval": "adjacent",
                        "retrievals": ["adjacent"],
                        "combined_score": 0.35,
                        "final_score": 0.35,
                    },
                )
            )

    return expanded[: max(FINAL_K * 2, 24)]


def build_context(docs: List[Document], max_chars: int) -> str:
    parts: List[str] = []
    used = 0

    for d in docs:
        source = d.metadata.get("source", "unknown")
        chunk = d.metadata.get("chunk", 0)
        retrievals = d.metadata.get("retrievals")
        retrieval = ",".join(retrievals) if retrievals else d.metadata.get("retrieval", "")
        part = f"[source={source} chunk={chunk} retrieval={retrieval}] {d.page_content}"
        part_len = len(part)

        if used + part_len > max_chars:
            remaining = max_chars - used
            if remaining > 250:
                parts.append(part[:remaining])
            break

        parts.append(part)
        used += part_len + 2

    return "\n\n".join(parts)


def rag_answer(workspace_id: str, question: str) -> str:
    query_variants = rewrite_queries(question)

    all_semantic: List[Document] = []
    all_keyword: List[Document] = []

    for q in query_variants:
        all_semantic.extend(semantic_search(workspace_id, q, k=SEMANTIC_K_PER_QUERY))
        all_keyword.extend(keyword_search(workspace_id, q, k=KEYWORD_K_PER_QUERY))

    docs = merge_results(all_semantic, all_keyword, question=question, final_k=FINAL_K)
    docs = expand_adjacent_docs(workspace_id, docs)
    context = build_context(docs, MAX_CONTEXT_CHARS)

    if not context.strip():
        return "I don't know based on the uploaded files."

    prompt = f"""{SYSTEM_PROMPT}

Context:
{context}

User question: {question}

Answer using only the context above. If the answer is partially supported, provide the strongest grounded answer possible and briefly say what remains unclear.
"""

    raw = generate_streaming(prompt)
    obj = safe_json_load(raw)
    answer = (obj.get("answer") or "").strip()

    if not answer:
        answer = "I don't know based on the uploaded files."

    if DEBUG_RETRIEVAL:
        st.session_state["last_debug_docs"] = [
            {
                "source": d.metadata.get("source"),
                "chunk": d.metadata.get("chunk"),
                "retrieval": d.metadata.get("retrieval"),
                "retrievals": d.metadata.get("retrievals", []),
                "preview": d.page_content[:400],
            }
            for d in docs
        ]

    touch_activity(workspace_id)
    return answer


# ============================
# STREAMLIT APP
# ============================
def main():
    st.set_page_config(page_title=APP_TITLE, layout="centered")
    inject_styles()

    ensure_dir(APP_DATA_ROOT)
    ensure_dir(CHROMA_ROOT)
    assert_writable_dir(APP_DATA_ROOT)
    assert_writable_dir(CHROMA_ROOT)

    janitor_sweep()
    login_gate()
    workspace_id = ensure_tenant_context()

    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "indexed" not in st.session_state:
        st.session_state.indexed = False
    if "failed_files" not in st.session_state:
        st.session_state.failed_files = []
    if "processed_files" not in st.session_state:
        st.session_state.processed_files = []
    if "created_at" not in st.session_state:
        st.session_state.created_at = _now()
    if "last_activity" not in st.session_state:
        st.session_state.last_activity = _now()
    if "last_debug_docs" not in st.session_state:
        st.session_state.last_debug_docs = []
    if "db_version" not in st.session_state:
        st.session_state.db_version = utc_stamp()

    manifest = load_manifest(workspace_id)
    if manifest.get("chunks"):
        st.session_state.indexed = True
        st.session_state.last_index_count = len(manifest.get("chunks", []))
        if not st.session_state.get("processed_files"):
            st.session_state.processed_files = [
                f.get("file_name", "") for f in manifest.get("files", [])
            ]

    touch_activity(workspace_id)

    st.markdown(f'<div class="gothic-title">{APP_TITLE}</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="center-note">Ask Questions About Your Documents</div>',
        unsafe_allow_html=True,
    )

    if not st.session_state.indexed:
        st.markdown(
            '<div class="gothic-sub">Upload Documents and Ask Questions Instantly</div>',
            unsafe_allow_html=True,
        )

        uploaded = st.file_uploader(
            "Upload files",
            type=None,
            accept_multiple_files=True,
            label_visibility="visible",
        )

        st.markdown(
            f'<div class="small-muted">Up to {MAX_FILES} files • {MAX_FILE_MB}MB max per file • {MAX_TOTAL_UPLOAD_MB}MB total</div>',
            unsafe_allow_html=True,
        )

        c1, c2 = st.columns([2, 1])
        with c1:
            index_clicked = st.button("Upload & Prepare", type="primary", use_container_width=True)
        with c2:
            if st.button("Reset", use_container_width=True):
                reset_workspace()
                st.rerun()

        if index_clicked:
            ok, msg = validate_uploads(uploaded)
            if not ok:
                st.error(msg)
            else:
                try:
                    with st.spinner("Preparing your files..."):
                        n = index_files(workspace_id, uploaded)
                except Exception as e:
                    st.error(
                        "Indexing failed due to a storage/database issue.\n\n"
                        f"Details: {e}"
                    )
                    return

                if n == 0:
                    st.warning(
                        "I couldn’t extract text from those files. Supported examples include PDF, DOC/DOCX, PPT/PPTX, ODT/ODS/ODP, TXT, CSV, XLS/XLSX/XLSB, Parquet, HTML, XML, YAML, JSON, JSONL, NDJSON, IPYNB, EML, MSG, RTF, RST, TEX, and code/text files."
                    )
                    if st.session_state.get("failed_files"):
                        st.error(
                            "Files that failed:\n\n"
                            + "\n".join(f"- {x}" for x in st.session_state["failed_files"])
                        )
                else:
                    st.success(
                        f"Ready. Indexed {n} chunks from {len(st.session_state.get('processed_files', []))} file(s)."
                    )
                    if st.session_state.get("failed_files"):
                        st.warning(
                            "Some files could not be processed:\n\n"
                            + "\n".join(f"- {x}" for x in st.session_state["failed_files"])
                        )
                    st.rerun()

        st.markdown(
            f"""
<div class="small-foot">
Files are automatically deleted after <strong>{IDLE_TTL_SECONDS // 60} minutes of inactivity</strong> or <strong>{WORKSPACE_TTL_SECONDS // 3600} hours maximum session time</strong>.<br/>
<em>This is a demo environment. Please upload sample documents only.</em><br/>
Questions or custom deployments: <strong>linkedin.com/in/thedannyscott</strong>
</div>
""",
            unsafe_allow_html=True,
        )
        return

    top = st.columns([3, 1])
    with top[0]:
        processed = st.session_state.get("processed_files", [])
        failed = st.session_state.get("failed_files", [])
        st.markdown(
            f'<div class="small-success">Indexed files: {len(processed)} • Failed files: {len(failed)} • Chunks: {st.session_state.get("last_index_count", 0)}</div>',
            unsafe_allow_html=True,
        )
    with top[1]:
        if st.button("Reset / New Upload", use_container_width=True):
            reset_workspace()
            st.rerun()

    if st.session_state.get("processed_files"):
        with st.expander("Show indexed file names"):
            st.write("\n".join(st.session_state["processed_files"]))

    if st.session_state.get("failed_files"):
        with st.expander("Show file processing issues"):
            st.write("\n".join(st.session_state["failed_files"]))

    if DEBUG_RETRIEVAL and st.session_state.get("last_debug_docs"):
        with st.expander("Debug: last retrieved chunks"):
            st.json(st.session_state["last_debug_docs"])

    for m in st.session_state.messages:
        avatar = ASSISTANT_AVATAR if m["role"] == "assistant" else None
        with st.chat_message(m["role"], avatar=avatar):
            st.markdown(m["content"])

    q = st.chat_input("Ask a question about your uploaded files…")
    if q:
        st.session_state.messages.append({"role": "user", "content": q})

        with st.chat_message("assistant", avatar=ASSISTANT_AVATAR):
            with st.spinner("Thinking..."):
                answer = rag_answer(workspace_id, q)
            st.markdown(answer)

        st.session_state.messages.append({"role": "assistant", "content": answer})


if __name__ == "__main__":
    main()