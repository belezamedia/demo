import os
import io
import re
import json
import uuid
import time
import hmac
import yaml
import hashlib
import zipfile
import configparser
import threading
from pathlib import Path
from typing import Any, Dict, List, Tuple
from urllib.parse import unquote, quote
from html.parser import HTMLParser
from email import policy
from email.parser import BytesParser
from datetime import datetime
from collections import Counter, defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
import xml.etree.ElementTree as ET

import streamlit as st
import pandas as pd
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain_aws import ChatBedrockConverse, BedrockEmbeddings

# Load .env locally only (never override Streamlit Cloud secrets)
try:
    from dotenv import load_dotenv

    if Path(".env").exists():
        load_dotenv(override=False)
except Exception:
    pass

# Optional parsers
try:
    from striprtf.striprtf import rtf_to_text as _rtf_to_text
except Exception:
    _rtf_to_text = None

try:
    import extract_msg  # .msg
except Exception:
    extract_msg = None

try:
    import olefile  # .doc / .ppt / .msg fallback
except Exception:
    olefile = None


# ============================
# CONFIG
# ============================
APP_TITLE = "DocHelp.AI"
BRAND_PINK = "#1F4ED8"
DEFAULT_REGION = os.getenv("AWS_DEFAULT_REGION", "us-west-2").strip()

BEDROCK_CHAT_MODEL_ID = unquote(
    os.getenv("BEDROCK_CHAT_MODEL_ID", "us.amazon.nova-pro-v1:0").strip()
)
BEDROCK_EMBED_MODEL_ID = unquote(
    os.getenv("BEDROCK_EMBED_MODEL_ID", "amazon.titan-embed-text-v2:0").strip()
)

TENANT_SIGNING_KEY = os.getenv("TENANT_SIGNING_KEY", "").strip()
APP_USERNAME = os.getenv("APP_USERNAME", "").strip()
APP_PASSWORD = os.getenv("APP_PASSWORD", "").strip()

MAX_FILES = int(os.getenv("MAX_FILES", "300"))
MAX_FILE_MB = int(os.getenv("MAX_FILE_MB", "500"))
MAX_TOTAL_UPLOAD_MB = int(os.getenv("MAX_TOTAL_UPLOAD_MB", "3000"))
MAX_TOTAL_CHUNKS = int(os.getenv("MAX_TOTAL_CHUNKS", "100000"))

SEMANTIC_K = int(os.getenv("SEMANTIC_K", "10"))
KEYWORD_K = int(os.getenv("KEYWORD_K", "10"))
FINAL_K = int(os.getenv("FINAL_K", "12"))
ADJACENT_CHUNKS = int(os.getenv("ADJACENT_CHUNKS", "1"))

CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "1200"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "120"))
MAX_CONTEXT_CHARS = int(os.getenv("MAX_CONTEXT_CHARS", "22000"))
MAX_PARSE_WORKERS = int(os.getenv("MAX_PARSE_WORKERS", "6"))
ENABLE_QUERY_REWRITE = os.getenv("ENABLE_QUERY_REWRITE", "false").strip().lower() == "true"
DEBUG_RETRIEVAL = os.getenv("DEBUG_RETRIEVAL", "false").strip().lower() == "true"

WORKSPACE_TTL_SECONDS = int(os.getenv("WORKSPACE_TTL_SECONDS", "7200"))
IDLE_TTL_SECONDS = int(os.getenv("IDLE_TTL_SECONDS", "1800"))

SYSTEM_PROMPT = """You are a professional business assistant.

Rules:
1. Answer using only the provided context from the uploaded files.
2. Do not use outside knowledge.
3. If the answer is partially supported by the uploaded files, provide the best grounded answer and clearly note what is uncertain or missing.
4. If the answer is not supported by the uploaded files, say exactly:
   \"I don't know based on the uploaded files.\"
5. Write in a concise, professional, businesslike tone.
6. Prefer a direct answer first.
7. When useful, mention the file names that support the answer.
8. Never invent policy details, numbers, dates, entities, or conclusions not present in the context.

Return VALID JSON only in this format:
{"answer":"..."}

No markdown. No extra keys. No explanations.
"""

QUERY_REWRITE_PROMPT = """You rewrite user questions into search-friendly variants for document retrieval.

Rules:
- Stay faithful to the original meaning.
- Produce 3 concise alternate queries.
- Include likely key terms, synonyms, and document-oriented phrasing.
- Do not add facts not present in the question.
- If the question mentions a file name, preserve it.
- Return valid JSON only:
{"queries":["q1","q2","q3"]}
"""

_STOPWORDS = {
    "the", "a", "an", "and", "or", "but", "if", "then", "of", "to", "in", "on", "for",
    "with", "by", "from", "at", "as", "is", "are", "was", "were", "be", "been", "being",
    "this", "that", "these", "those", "it", "its", "their", "there", "here", "what",
    "which", "who", "whom", "how", "when", "where", "why", "can", "could", "should",
    "would", "do", "does", "did", "about", "into", "than", "them", "they", "you", "your",
    "me", "my", "we", "our", "please", "show", "tell", "give"
}

_PARSE_LOCKS: Dict[str, threading.Lock] = {}


# ============================
# UI / STYLES
# ============================
def inject_styles() -> None:
    st.markdown(
        f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

:root {{
  --brand-blue: {BRAND_PINK};
}}

html, body, [class*="css"] {{
  font-family: 'Inter', sans-serif;
}}

.gothic-title {{
  font-family: 'Inter', sans-serif;
  font-size: 56px;
  font-weight: 700;
  line-height: 1.0;
  text-align: center;
  margin-top: 0.5rem;
  margin-bottom: 0.35rem;
}}

.gothic-sub {{
  font-family: 'Inter', sans-serif;
  font-size: 20px;
  font-weight: 600;
  text-align: center;
  margin-bottom: 0.75rem;
}}

.center-note {{
  text-align: center;
  font-size: 18px;
  color: rgba(0,0,0,0.72);
  margin-bottom: 1rem;
}}

.small-foot {{
  text-align: center;
  font-size: 13px;
  color: rgba(0,0,0,0.65);
  margin-top: 1rem;
  line-height: 1.65;
}}

.small-foot em {{
  font-style: italic;
}}

.small-muted {{
  text-align: center;
  font-size: 13px;
  color: rgba(0,0,0,0.58);
  margin-top: 0.25rem;
  margin-bottom: 0.5rem;
}}

.small-success {{
  font-size: 13px;
  color: rgba(0,0,0,0.72);
  margin-top: 0.25rem;
  margin-bottom: 0.5rem;
}}

#MainMenu {{visibility: hidden;}}
footer {{visibility: hidden;}}
header {{visibility: hidden;}}

button[kind="primary"] {{
  background: var(--brand-blue) !important;
  border: 1px solid var(--brand-blue) !important;
  color: white !important;
  font-family: 'Inter', sans-serif !important;
}}

div.stButton > button {{
  border-radius: 12px !important;
  padding: 0.65rem 1rem !important;
}}

div[data-testid="stTextInput"] input {{
  border: 2px solid var(--brand-blue) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
}}

div[data-testid="stTextInput"] input:focus {{
  border: 2px solid var(--brand-blue) !important;
  outline: none !important;
  box-shadow: 0 0 0 3px rgba(31, 78, 216, 0.15) !important;
}}

div[data-testid="stChatInput"] textarea {{
  border: 2px solid var(--brand-blue) !important;
  border-radius: 14px !important;
  box-shadow: none !important;
}}

div[data-testid="stChatInput"] textarea:focus {{
  border: 2px solid var(--brand-blue) !important;
  outline: none !important;
  box-shadow: 0 0 0 3px rgba(31, 78, 216, 0.15) !important;
}}

div[data-testid="stChatInput"] button {{
  border-radius: 12px !important;
  border: 1px solid var(--brand-blue) !important;
  background: var(--brand-blue) !important;
}}

div[data-testid="stChatInput"] button svg {{
  fill: white !important;
  stroke: white !important;
}}
</style>
""",
        unsafe_allow_html=True,
    )


# ============================
# LOGIN GATE
# ============================
def login_gate() -> None:
    if st.session_state.get("authed") is True:
        return

    st.markdown(f'<div class="gothic-title">{APP_TITLE}</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="center-note">Private AI for Your Documents</div>',
        unsafe_allow_html=True,
    )

    if not (APP_USERNAME and APP_PASSWORD):
        st.error(
            "Missing APP_USERNAME / APP_PASSWORD. Set them in .env (local) or Streamlit Secrets (cloud)."
        )
        st.stop()

    st.markdown('<div class="gothic-sub">Login</div>', unsafe_allow_html=True)

    u = st.text_input("username", label_visibility="collapsed", placeholder="username")
    p = st.text_input(
        "password",
        type="password",
        label_visibility="collapsed",
        placeholder="password",
    )

    c1, c2 = st.columns([2, 1])
    with c1:
        go = st.button("Enter", type="primary", use_container_width=True)
    with c2:
        if st.button("Reset", use_container_width=True):
            st.session_state.clear()
            st.rerun()

    if go:
        if u == APP_USERNAME and p == APP_PASSWORD:
            st.session_state["authed"] = True
            st.rerun()
        else:
            st.error("Invalid username or password.")

    st.stop()


# ============================
# TENANT ISOLATION
# ============================
def _b(x: str) -> bytes:
    return (x or "").encode("utf-8")


def _is_streamlit_cloud() -> bool:
    return Path("/mount/src").exists()


def sign_token(token: str) -> str:
    if _is_streamlit_cloud() and not TENANT_SIGNING_KEY:
        raise RuntimeError("TENANT_SIGNING_KEY must be set in Streamlit Secrets.")
    key = TENANT_SIGNING_KEY or "insecure-dev-key-change-me"
    return hmac.new(_b(key), _b(token), hashlib.sha256).hexdigest()


def verify_token(token: str, sig: str) -> bool:
    if not token or not sig:
        return False
    expected = sign_token(token)
    return hmac.compare_digest(expected, sig)


def workspace_id_from_token(token: str) -> str:
    return hashlib.sha256(_b(token)).hexdigest()


def set_query_params(token: str, sig: str) -> None:
    try:
        st.query_params["t"] = token
        st.query_params["sig"] = sig
    except Exception:
        try:
            st.experimental_set_query_params(t=token, sig=sig)
        except Exception:
            pass


def get_query_params() -> Dict[str, str]:
    try:
        qp = st.query_params
        return {"t": qp.get("t", ""), "sig": qp.get("sig", "")}
    except Exception:
        qp = st.experimental_get_query_params()
        return {
            "t": (qp.get("t", [""]) or [""])[0],
            "sig": (qp.get("sig", [""]) or [""])[0],
        }


def ensure_tenant_context() -> str:
    qp = get_query_params()
    token = (qp.get("t") or "").strip()
    sig = (qp.get("sig") or "").strip()

    valid = False
    try:
        valid = verify_token(token, sig)
    except Exception:
        valid = False

    if not valid:
        token = uuid.uuid4().hex
        sig = sign_token(token)
        set_query_params(token, sig)

    workspace_id = workspace_id_from_token(token)
    st.session_state["workspace_id"] = workspace_id
    return workspace_id


# ============================
# AVATAR
# ============================
def assistant_avatar_data_uri(color_hex: str) -> str:
    svg = f"""
    <svg xmlns="http://www.w3.org/2000/svg" width="64" height="64" viewBox="0 0 64 64">
      <rect x="14" y="18" width="36" height="32" rx="10" fill="{color_hex}"/>
      <rect x="26" y="10" width="12" height="10" rx="4" fill="{color_hex}"/>
      <circle cx="26" cy="34" r="4" fill="white"/>
      <circle cx="38" cy="34" r="4" fill="white"/>
      <rect x="24" y="42" width="16" height="4" rx="2" fill="white" opacity="0.95"/>
      <circle cx="32" cy="8" r="3" fill="{color_hex}"/>
      <rect x="31" y="8" width="2" height="6" fill="{color_hex}"/>
    </svg>
    """.strip()
    return "data:image/svg+xml;utf8," + quote(svg)


ASSISTANT_AVATAR = assistant_avatar_data_uri(BRAND_PINK)


# ============================
# CACHED RESOURCES
# ============================
@st.cache_resource(show_spinner=False)
def get_splitter() -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", "; ", ", ", " ", ""],
    )


@st.cache_resource(show_spinner=False)
def get_embeddings() -> BedrockEmbeddings:
    return BedrockEmbeddings(
        region_name=DEFAULT_REGION,
        model_id=BEDROCK_EMBED_MODEL_ID,
    )


@st.cache_resource(show_spinner=False)
def get_llm() -> ChatBedrockConverse:
    return ChatBedrockConverse(
        model=BEDROCK_CHAT_MODEL_ID,
        region_name=DEFAULT_REGION,
        temperature=0.0,
        max_tokens=700,
    )


# ============================
# SESSION INDEX STORAGE
# ============================
def _now() -> int:
    return int(time.time())


def utc_stamp() -> str:
    return datetime.utcnow().strftime("%Y%m%d%H%M%S%f")


def ensure_workspace_state(workspace_id: str) -> Dict[str, Any]:
    key = f"index_data::{workspace_id}"
    if key not in st.session_state:
        st.session_state[key] = {
            "files": [],
            "chunks": [],
            "chunk_lookup": {},
            "inverted_index": defaultdict(set),
            "file_hashes": set(),
            "created_at": _now(),
            "last_activity": _now(),
        }
    return st.session_state[key]


def touch_activity(workspace_id: str) -> None:
    data = ensure_workspace_state(workspace_id)
    data["last_activity"] = _now()
    st.session_state["last_activity"] = data["last_activity"]
    st.session_state.setdefault("created_at", data.get("created_at", _now()))


def workspace_expired(workspace_id: str) -> bool:
    data = ensure_workspace_state(workspace_id)
    created_at = int(data.get("created_at") or _now())
    last_activity = int(data.get("last_activity") or _now())
    now = _now()
    return (now - created_at) > WORKSPACE_TTL_SECONDS or (now - last_activity) > IDLE_TTL_SECONDS


def clear_workspace_storage(workspace_id: str) -> None:
    ensure_workspace_state(workspace_id)
    st.session_state[f"index_data::{workspace_id}"] = {
        "files": [],
        "chunks": [],
        "chunk_lookup": {},
        "inverted_index": defaultdict(set),
        "file_hashes": set(),
        "created_at": _now(),
        "last_activity": _now(),
    }
    st.session_state[f"vectordb::{workspace_id}"] = Chroma(
        collection_name=f"demo_docs_{workspace_id[:12]}_{utc_stamp()}",
        embedding_function=get_embeddings(),
    )


def get_index_data(workspace_id: str) -> Dict[str, Any]:
    return ensure_workspace_state(workspace_id)


def get_vectordb(workspace_id: str) -> Chroma:
    key = f"vectordb::{workspace_id}"
    if key not in st.session_state:
        st.session_state[key] = Chroma(
            collection_name=f"demo_docs_{workspace_id[:12]}",
            embedding_function=get_embeddings(),
        )
    return st.session_state[key]


def reset_workspace() -> None:
    token = uuid.uuid4().hex
    sig = sign_token(token)
    set_query_params(token, sig)

    st.session_state.indexed = False
    st.session_state.last_index_count = 0
    st.session_state.messages = []
    st.session_state.failed_files = []
    st.session_state.processed_files = []
    st.session_state.last_debug_docs = []
    st.session_state.created_at = _now()
    st.session_state.last_activity = _now()


# ============================
# BEDROCK STREAMING
# ============================
def _extract_text(obj: Any) -> str:
    if isinstance(obj, str):
        return obj
    if hasattr(obj, "content"):
        return _extract_text(obj.content)
    if isinstance(obj, list):
        parts: List[str] = []
        for x in obj:
            if isinstance(x, dict):
                parts.append(str(x.get("text", "")))
            elif hasattr(x, "get"):
                parts.append(str(x.get("text", "")))
            else:
                parts.append(str(x))
        return "".join(parts)
    return ""


def generate_streaming(prompt: str) -> str:
    llm = get_llm()
    parts: List[str] = []
    try:
        for chunk in llm.stream(prompt):
            txt = _extract_text(chunk)
            if txt:
                parts.append(txt)
        return "".join(parts)
    except Exception:
        resp = llm.invoke(prompt)
        return _extract_text(resp)


def safe_json_load(s: str) -> Dict[str, Any]:
    try:
        obj = json.loads(s)
        if isinstance(obj, dict):
            return obj
    except Exception:
        pass

    m = re.search(r"\{.*\}", s, flags=re.DOTALL)
    if m:
        try:
            obj = json.loads(m.group(0))
            if isinstance(obj, dict):
                return obj
        except Exception:
            pass

    return {"answer": s.strip()}


def dedupe_preserve_order(items: List[str]) -> List[str]:
    out: List[str] = []
    seen = set()
    for item in items:
        if item not in seen:
            seen.add(item)
            out.append(item)
    return out


def rewrite_queries(question: str) -> List[str]:
    if not ENABLE_QUERY_REWRITE:
        return [question]

    prompt = f"""{QUERY_REWRITE_PROMPT}

Original question:
{question}
"""
    try:
        raw = generate_streaming(prompt)
        obj = safe_json_load(raw)
        queries = obj.get("queries", [])
        if isinstance(queries, list):
            clean: List[str] = []
            for q in queries:
                q = str(q or "").strip()
                if q and q.lower() != question.lower():
                    clean.append(q)
            return dedupe_preserve_order([question] + clean[:3])
    except Exception:
        pass
    return [question]


# ============================
# FILE PARSING HELPERS
# ============================
class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self._chunks: List[str] = []

    def handle_data(self, data: str) -> None:
        if data and data.strip():
            self._chunks.append(data.strip())

    def get_text(self) -> str:
        return " ".join(self._chunks).strip()


def _clean_whitespace(s: str) -> str:
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\r\n?", "\n", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _decode_best_effort(file_bytes: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return file_bytes.decode(enc)
        except Exception:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def _strip_xml_tags(xml_text: str) -> str:
    try:
        root = ET.fromstring(xml_text)
        parts: List[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        return _clean_whitespace("\n".join(parts))
    except Exception:
        return _clean_whitespace(re.sub(r"<[^>]+>", " ", xml_text))


def _read_zip_member_text(file_bytes: bytes, member_names: List[str]) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            parts: List[str] = []
            names = set(zf.namelist())
            for name in member_names:
                if name in names:
                    xml_text = zf.read(name).decode("utf-8", errors="ignore")
                    cleaned = _strip_xml_tags(xml_text)
                    if cleaned:
                        parts.append(cleaned)
            return _clean_whitespace("\n\n".join(parts))
    except Exception:
        return ""


def read_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    texts: List[str] = []
    for i, page in enumerate(reader.pages):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if txt.strip():
            texts.append(f"[Page {i + 1}]\n{txt}")
    return _clean_whitespace("\n\n".join(texts))


def read_csv(file_bytes: bytes, sep: str = ",") -> str:
    last_err = None
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            df = pd.read_csv(io.BytesIO(file_bytes), sep=sep, encoding=enc)
            return df.to_csv(index=False)
        except Exception as e:
            last_err = e
    raise RuntimeError(f"Could not read delimited file: {last_err}")


def _dataframe_to_text(df: pd.DataFrame) -> str:
    if df is None:
        return ""
    if df.empty:
        cols = [str(c) for c in df.columns]
        return ",".join(cols) if cols else ""

    rows = min(len(df), 5000)
    cols = min(len(df.columns), 100)
    slim = df.iloc[:rows, :cols].copy()
    return slim.fillna("").to_csv(index=False)


def read_excel(file_bytes: bytes, filename: str = "") -> str:
    name = (filename or "").lower()
    bio = io.BytesIO(file_bytes)
    try:
        if name.endswith(".xlsx"):
            sheets = pd.read_excel(bio, engine="openpyxl", sheet_name=None)
        elif name.endswith(".xls"):
            sheets = pd.read_excel(bio, engine="xlrd", sheet_name=None)
        else:
            sheets = pd.read_excel(bio, sheet_name=None)
    except ImportError as e:
        raise RuntimeError(
            f"Excel dependency missing for '{filename}'. Make sure openpyxl is installed for .xlsx and xlrd for .xls."
        ) from e
    except Exception as e:
        raise RuntimeError(f"Could not read Excel file '{filename}': {e}") from e

    parts: List[str] = []
    for sheet_name, df in (sheets or {}).items():
        parts.append(f"=== Sheet: {sheet_name} ===")
        parts.append(_dataframe_to_text(df))
        parts.append("")
    return _clean_whitespace("\n".join(parts))


def read_xlsb(file_bytes: bytes, filename: str = "") -> str:
    try:
        sheets = pd.read_excel(io.BytesIO(file_bytes), engine="pyxlsb", sheet_name=None)
        parts: List[str] = []
        for sheet_name, df in (sheets or {}).items():
            parts.append(f"=== Sheet: {sheet_name} ===")
            parts.append(_dataframe_to_text(df))
            parts.append("")
        return _clean_whitespace("\n".join(parts))
    except ImportError as e:
        raise RuntimeError(
            f"Excel dependency missing for '{filename}'. Make sure pyxlsb is installed for .xlsb files."
        ) from e
    except Exception as e:
        raise RuntimeError(f"Could not read XLSB file '{filename}': {e}") from e


def read_parquet(file_bytes: bytes, filename: str = "") -> str:
    try:
        df = pd.read_parquet(io.BytesIO(file_bytes))
        return _dataframe_to_text(df)
    except Exception as e:
        raise RuntimeError(f"Could not read Parquet file '{filename}': {e}") from e


def read_json(file_bytes: bytes) -> str:
    try:
        obj = json.loads(_decode_best_effort(file_bytes))
        return json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception:
        return _decode_best_effort(file_bytes).strip()


def read_jsonl(file_bytes: bytes) -> str:
    text = _decode_best_effort(file_bytes)
    rows: List[str] = []
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            obj = json.loads(line)
            rows.append(json.dumps(obj, ensure_ascii=False))
        except Exception:
            rows.append(line)
    return _clean_whitespace("\n".join(rows))


def read_yaml(file_bytes: bytes) -> str:
    text = _decode_best_effort(file_bytes)
    try:
        obj = yaml.safe_load(text)
        return yaml.safe_dump(obj, sort_keys=False, allow_unicode=True)
    except Exception:
        return text.strip()


def read_ini(file_bytes: bytes) -> str:
    raw = _decode_best_effort(file_bytes)
    cp = configparser.ConfigParser()
    try:
        cp.read_string(raw)
        out_lines: List[str] = []
        for section in cp.sections():
            out_lines.append(f"[{section}]")
            for k, v in cp.items(section):
                out_lines.append(f"{k} = {v}")
            out_lines.append("")
        return _clean_whitespace("\n".join(out_lines))
    except Exception:
        return raw.strip()


def read_xml(file_bytes: bytes) -> str:
    try:
        root = ET.fromstring(file_bytes)
        parts: List[str] = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                parts.append(elem.tail.strip())
        return _clean_whitespace("\n".join(parts))
    except Exception:
        xml_text = _decode_best_effort(file_bytes)
        return _clean_whitespace(re.sub(r"<[^>]+>", " ", xml_text))


def read_html(file_bytes: bytes) -> str:
    html = _decode_best_effort(file_bytes)
    parser = _HTMLTextExtractor()
    try:
        parser.feed(html)
        return _clean_whitespace(parser.get_text())
    except Exception:
        return _clean_whitespace(re.sub(r"<[^>]+>", " ", html))


def read_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return _clean_whitespace("\n".join(parts))


def read_pptx(file_bytes: bytes) -> str:
    pres = Presentation(io.BytesIO(file_bytes))
    parts: List[str] = []
    for slide_idx, slide in enumerate(pres.slides, start=1):
        slide_parts: List[str] = []
        for shape in slide.shapes:
            txt = getattr(shape, "text", None)
            if txt and str(txt).strip():
                slide_parts.append(str(txt).strip())
        if slide_parts:
            parts.append(f"[Slide {slide_idx}]")
            parts.extend(slide_parts)
            parts.append("")
    return _clean_whitespace("\n".join(parts))


def read_odt(file_bytes: bytes) -> str:
    return _read_zip_member_text(file_bytes, ["content.xml", "styles.xml", "meta.xml"])


def read_ods(file_bytes: bytes) -> str:
    return _read_zip_member_text(file_bytes, ["content.xml", "styles.xml", "meta.xml"])


def read_odp(file_bytes: bytes) -> str:
    return _read_zip_member_text(file_bytes, ["content.xml", "styles.xml", "meta.xml"])


def _extract_printable_utf16le(blob: bytes) -> List[str]:
    text = blob.decode("utf-16le", errors="ignore")
    candidates = re.findall(r"[\x20-\x7E][\x20-\x7E\n\r\t]{3,}", text)
    return [c.strip() for c in candidates if c.strip()]


def _extract_printable_utf8(blob: bytes) -> List[str]:
    text = blob.decode("utf-8", errors="ignore")
    candidates = re.findall(r"[A-Za-z0-9][^\x00]{3,}", text)
    return [c.strip() for c in candidates if c.strip()]


def read_doc_legacy(file_bytes: bytes) -> str:
    if olefile is None:
        return ""
    try:
        with olefile.OleFileIO(io.BytesIO(file_bytes)) as ole:
            pieces: List[str] = []
            for stream_name in ole.listdir():
                joined = "/".join(stream_name)
                if joined in ("WordDocument", "1Table", "0Table"):
                    try:
                        blob = ole.openstream(stream_name).read()
                        pieces.extend(_extract_printable_utf16le(blob))
                        pieces.extend(_extract_printable_utf8(blob))
                    except Exception:
                        continue
            return _clean_whitespace("\n".join(dict.fromkeys(pieces)))
    except Exception:
        return ""


def read_ppt_legacy(file_bytes: bytes) -> str:
    if olefile is None:
        return ""
    try:
        with olefile.OleFileIO(io.BytesIO(file_bytes)) as ole:
            pieces: List[str] = []
            for stream_name in ole.listdir():
                joined = "/".join(stream_name)
                if joined == "PowerPoint Document":
                    try:
                        blob = ole.openstream(stream_name).read()
                        pieces.extend(_extract_printable_utf16le(blob))
                        pieces.extend(_extract_printable_utf8(blob))
                    except Exception:
                        continue
            return _clean_whitespace("\n".join(dict.fromkeys(pieces)))
    except Exception:
        return ""


def read_msg(file_bytes: bytes) -> str:
    if extract_msg is not None:
        tmp_path = None
        try:
            tmp_dir = Path("/tmp")
            tmp_dir.mkdir(parents=True, exist_ok=True)
            tmp_path = tmp_dir / f"msg_{uuid.uuid4().hex}.msg"
            tmp_path.write_bytes(file_bytes)

            msg = extract_msg.Message(str(tmp_path))
            parts: List[str] = []
            if getattr(msg, "subject", None):
                parts.append(f"Subject: {msg.subject}")
            if getattr(msg, "sender", None):
                parts.append(f"From: {msg.sender}")
            if getattr(msg, "to", None):
                parts.append(f"To: {msg.to}")
            if getattr(msg, "date", None):
                parts.append(f"Date: {msg.date}")
            body = getattr(msg, "body", "") or ""
            if body.strip():
                parts.append("")
                parts.append(body.strip())
            return _clean_whitespace("\n".join(parts))
        except Exception:
            pass
        finally:
            try:
                if tmp_path and tmp_path.exists():
                    tmp_path.unlink()
            except Exception:
                pass

    if olefile is None:
        return ""

    try:
        with olefile.OleFileIO(io.BytesIO(file_bytes)) as ole:
            parts: List[str] = []
            for stream_name in ole.listdir():
                joined = "/".join(stream_name)
                if "__substg1.0_" in joined:
                    try:
                        blob = ole.openstream(stream_name).read()
                        parts.extend(_extract_printable_utf16le(blob))
                    except Exception:
                        continue
            return _clean_whitespace("\n".join(dict.fromkeys(parts)))
    except Exception:
        return ""


def read_ipynb(file_bytes: bytes) -> str:
    try:
        nb = json.loads(_decode_best_effort(file_bytes))
        cells = nb.get("cells", []) if isinstance(nb, dict) else []
        parts: List[str] = []
        for cell in cells:
            src = cell.get("source", [])
            s = "".join(src) if isinstance(src, list) else str(src or "")
            if s.strip():
                parts.append(s.strip())
        return _clean_whitespace("\n\n".join(parts))
    except Exception:
        return _decode_best_effort(file_bytes).strip()


def read_eml(file_bytes: bytes) -> str:
    try:
        msg = BytesParser(policy=policy.default).parsebytes(file_bytes)
        parts: List[str] = []
        subj = msg.get("subject", "")
        frm = msg.get("from", "")
        to = msg.get("to", "")
        if subj:
            parts.append(f"Subject: {subj}")
        if frm:
            parts.append(f"From: {frm}")
        if to:
            parts.append(f"To: {to}")
        parts.append("")

        body_texts: List[str] = []
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                disp = part.get_content_disposition()
                if disp == "attachment":
                    continue
                if ctype in ("text/plain", "text/html"):
                    try:
                        content = part.get_content()
                    except Exception:
                        content = part.get_payload(decode=True)
                        content = (content or b"").decode("utf-8", errors="ignore")
                    if ctype == "text/html":
                        content = read_html(str(content).encode("utf-8", errors="ignore"))
                    if content and str(content).strip():
                        body_texts.append(str(content).strip())
        else:
            ctype = msg.get_content_type()
            content = msg.get_content()
            if ctype == "text/html":
                content = read_html(str(content).encode("utf-8", errors="ignore"))
            if content and str(content).strip():
                body_texts.append(str(content).strip())

        if body_texts:
            parts.append("\n\n".join(body_texts))

        return _clean_whitespace("\n".join(parts))
    except Exception:
        return _decode_best_effort(file_bytes).strip()


def read_rtf(file_bytes: bytes) -> str:
    raw = _decode_best_effort(file_bytes)
    if _rtf_to_text is not None:
        try:
            return _clean_whitespace(_rtf_to_text(raw))
        except Exception:
            pass
    fallback = re.sub(r"{\\.*?}|\\[a-zA-Z]+\d* ?", " ", raw)
    return _clean_whitespace(fallback)


def read_txt(file_bytes: bytes) -> str:
    return _clean_whitespace(_decode_best_effort(file_bytes))


CODE_EXTS = (
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cpp", ".h", ".hpp", ".go",
    ".rs", ".rb", ".php", ".cs", ".swift", ".kt", ".scala", ".sh", ".bash", ".zsh",
    ".ps1", ".sql", ".toml", ".env", ".dockerfile"
)
TEXTY_EXTS = (
    ".txt", ".md", ".log", ".cfg", ".conf", ".ini", ".rst", ".tex"
)


def file_to_text(uploaded_file) -> str:
    name = (uploaded_file.name or "").lower()
    b = uploaded_file.getvalue()

    if name.endswith(".pdf"):
        return read_pdf(b)
    if name.endswith(".docx"):
        return read_docx(b)
    if name.endswith(".doc"):
        return read_doc_legacy(b)
    if name.endswith(".pptx"):
        return read_pptx(b)
    if name.endswith(".ppt"):
        return read_ppt_legacy(b)
    if name.endswith(".odt"):
        return read_odt(b)
    if name.endswith(".ods"):
        return read_ods(b)
    if name.endswith(".odp"):
        return read_odp(b)
    if name.endswith(".csv"):
        return read_csv(b, sep=",")
    if name.endswith(".tsv"):
        return read_csv(b, sep="\t")
    if name.endswith(".psv"):
        return read_csv(b, sep="|")
    if name.endswith(".xlsb"):
        return read_xlsb(b, name)
    if name.endswith((".xls", ".xlsx")):
        return read_excel(b, name)
    if name.endswith(".parquet"):
        return read_parquet(b, name)
    if name.endswith((".html", ".htm")):
        return read_html(b)
    if name.endswith(".xml"):
        return read_xml(b)
    if name.endswith((".yaml", ".yml")):
        return read_yaml(b)
    if name.endswith((".ini", ".cfg", ".conf")):
        return read_ini(b)
    if name.endswith(".ipynb"):
        return read_ipynb(b)
    if name.endswith(".eml"):
        return read_eml(b)
    if name.endswith(".msg"):
        return read_msg(b)
    if name.endswith(".rtf"):
        return read_rtf(b)
    if name.endswith(".json"):
        return read_json(b)
    if name.endswith((".jsonl", ".ndjson")):
        return read_jsonl(b)
    if name.endswith(CODE_EXTS) or name.endswith(TEXTY_EXTS):
        return read_txt(b)
    return ""


@st.cache_data(show_spinner=False, ttl=3600, max_entries=512)
def parse_file_cached(file_name: str, file_bytes: bytes) -> str:
    class UploadedShim:
        def __init__(self, name: str, data: bytes):
            self.name = name
            self._data = data

        def getvalue(self):
            return self._data

    shim = UploadedShim(file_name, file_bytes)
    return file_to_text(shim)


# ============================
# VALIDATION / INDEX HELPERS
# ============================
def validate_uploads(files) -> Tuple[bool, str]:
    if not files:
        return False, "Please upload at least one file."

    if len(files) > MAX_FILES:
        return False, f"Too many files. Max is {MAX_FILES}."

    total_bytes = 0
    for f in files:
        size_bytes = int(getattr(f, "size", 0) or 0)
        total_bytes += size_bytes
        size_mb = size_bytes / (1024 * 1024)
        if size_mb > MAX_FILE_MB:
            return False, f"File '{f.name}' is too large ({size_mb:.1f} MB). Max is {MAX_FILE_MB} MB."

    total_mb = total_bytes / (1024 * 1024)
    if total_mb > MAX_TOTAL_UPLOAD_MB:
        return False, f"Total upload too large ({total_mb:.1f} MB). Max total is {MAX_TOTAL_UPLOAD_MB} MB."

    return True, ""


def _normalize_text_for_search(s: str) -> str:
    s = (s or "").lower()
    s = re.sub(r"[^a-z0-9\s._/-]", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def tokenize(s: str) -> List[str]:
    s = _normalize_text_for_search(s)
    return [t for t in s.split() if t and t not in _STOPWORDS and len(t) > 1]


def file_hash_bytes(file_bytes: bytes) -> str:
    return hashlib.sha256(file_bytes).hexdigest()


def build_chunk_id(file_hash: str, chunk_idx: int) -> str:
    return f"{file_hash}::{chunk_idx}"


def _parse_one_file(uploaded_file) -> Dict[str, Any]:
    raw_bytes = uploaded_file.getvalue()
    file_hash = file_hash_bytes(raw_bytes)
    text = parse_file_cached(uploaded_file.name, raw_bytes)
    return {
        "name": uploaded_file.name,
        "size_bytes": int(getattr(uploaded_file, "size", 0) or 0),
        "file_hash": file_hash,
        "text": text,
    }


def _batched(seq: List[Any], batch_size: int) -> List[List[Any]]:
    return [seq[i:i + batch_size] for i in range(0, len(seq), batch_size)]


def index_files(workspace_id: str, files) -> int:
    clear_workspace_storage(workspace_id)

    splitter = get_splitter()
    index_data = get_index_data(workspace_id)

    docs: List[Document] = []
    ids: List[str] = []
    failed_files: List[str] = []
    processed_files: List[str] = []
    parsed_results: List[Dict[str, Any]] = []

    with ThreadPoolExecutor(max_workers=max(1, MAX_PARSE_WORKERS)) as executor:
        futures = {executor.submit(_parse_one_file, f): f for f in files}
        for future in as_completed(futures):
            file_obj = futures[future]
            try:
                parsed_results.append(future.result())
            except Exception as e:
                failed_files.append(f"{file_obj.name}: {type(e).__name__}: {e}")

    parsed_results.sort(key=lambda x: x["name"].lower())
    total_chunks = 0

    for item in parsed_results:
        file_name = item["name"]
        file_hash = item["file_hash"]
        text = (item["text"] or "").strip()

        if not text:
            failed_files.append(f"{file_name}: no extractable text found")
            continue

        try:
            chunks = splitter.split_text(text)
        except Exception as e:
            failed_files.append(f"{file_name}: chunking failed: {e}")
            continue

        clean_chunks = [chunk.strip() for chunk in chunks if chunk and chunk.strip()]
        if not clean_chunks:
            failed_files.append(f"{file_name}: no usable chunks after parsing")
            continue

        if total_chunks + len(clean_chunks) > MAX_TOTAL_CHUNKS:
            failed_files.append("Chunk limit reached; some later file content may not have been indexed.")
            break

        index_data["files"].append(
            {
                "file_name": file_name,
                "file_hash": file_hash,
                "size_bytes": item["size_bytes"],
                "chunk_count": len(clean_chunks),
                "indexed_at": _now(),
            }
        )

        for idx, chunk in enumerate(clean_chunks):
            chunk_id = build_chunk_id(file_hash, idx)
            record = {
                "id": chunk_id,
                "source": file_name,
                "file_hash": file_hash,
                "chunk": idx,
                "text": chunk,
                "norm_text": _normalize_text_for_search(chunk),
            }
            index_data["chunks"].append(record)
            index_data["chunk_lookup"][(file_hash, idx)] = record

            for tok in set(tokenize(chunk) + tokenize(file_name)):
                index_data["inverted_index"][tok].add(chunk_id)

            docs.append(
                Document(
                    page_content=chunk,
                    metadata={
                        "id": chunk_id,
                        "source": file_name,
                        "file_hash": file_hash,
                        "chunk": idx,
                    },
                )
            )
            ids.append(chunk_id)
            total_chunks += 1

        processed_files.append(file_name)
        index_data["file_hashes"].add(file_hash)

    if docs:
        vectordb = get_vectordb(workspace_id)
        batch_size = 128
        for doc_batch, id_batch in zip(_batched(docs, batch_size), _batched(ids, batch_size)):
            vectordb.add_documents(doc_batch, ids=id_batch)

    st.session_state.indexed = len(index_data["chunks"]) > 0
    st.session_state.last_index_count = len(index_data["chunks"])
    st.session_state.failed_files = dedupe_preserve_order(failed_files)
    st.session_state.processed_files = dedupe_preserve_order(processed_files)
    touch_activity(workspace_id)
    return len(index_data["chunks"])


# ============================
# RETRIEVAL
# ============================
def score_keyword_match(question: str, chunk_text: str, filename: str = "") -> float:
    q_tokens = tokenize(question)
    if not q_tokens:
        return 0.0

    text_norm = _normalize_text_for_search(chunk_text)
    file_norm = _normalize_text_for_search(filename)
    text_tokens = text_norm.split()

    if not text_tokens and not file_norm:
        return 0.0

    q_counter = Counter(q_tokens)
    t_counter = Counter(text_tokens)

    overlap = 0.0
    for tok, q_count in q_counter.items():
        overlap += min(q_count, t_counter.get(tok, 0))

    phrase_bonus = 0.0
    q_norm = _normalize_text_for_search(question)
    if q_norm and q_norm in text_norm:
        phrase_bonus += 5.0

    exact_token_bonus = 0.0
    filename_bonus = 0.0
    for tok in set(q_tokens):
        if tok in text_norm:
            exact_token_bonus += 0.25
        if tok in file_norm:
            filename_bonus += 0.80

    density = overlap / max(1.0, len(q_tokens))
    return float(density * 10.0 + phrase_bonus + exact_token_bonus + filename_bonus)


def keyword_search(workspace_id: str, question: str, k: int = KEYWORD_K) -> List[Document]:
    index_data = get_index_data(workspace_id)
    chunks = index_data.get("chunks", [])
    inverted_index = index_data.get("inverted_index", {})

    q_tokens = tokenize(question)
    if not chunks or not q_tokens:
        return []

    candidate_ids = set()
    for tok in q_tokens:
        candidate_ids.update(inverted_index.get(tok, set()))

    if not candidate_ids:
        return []

    chunk_map = {r["id"]: r for r in chunks}
    scored: List[Tuple[float, Dict[str, Any]]] = []

    for cid in candidate_ids:
        record = chunk_map.get(cid)
        if not record:
            continue
        score = score_keyword_match(question, record.get("text", ""), record.get("source", ""))
        if score > 0:
            scored.append((score, record))

    scored.sort(key=lambda x: x[0], reverse=True)

    docs: List[Document] = []
    for score, record in scored[:k]:
        docs.append(
            Document(
                page_content=record["text"],
                metadata={
                    "id": record["id"],
                    "source": record["source"],
                    "file_hash": record["file_hash"],
                    "chunk": record["chunk"],
                    "retrieval": "keyword",
                    "keyword_score": round(score, 4),
                },
            )
        )
    return docs


def semantic_search(workspace_id: str, question: str, k: int = SEMANTIC_K) -> List[Document]:
    vectordb = get_vectordb(workspace_id)
    try:
        docs = vectordb.max_marginal_relevance_search(question, k=k, fetch_k=max(k * 2, 20))
    except Exception:
        try:
            docs = vectordb.similarity_search(question, k=k)
        except Exception:
            docs = []

    out: List[Document] = []
    for rank, doc in enumerate(docs):
        md = dict(doc.metadata or {})
        md["retrieval"] = "semantic"
        md["semantic_score"] = max(0.0, float(k - rank))
        out.append(Document(page_content=doc.page_content, metadata=md))
    return out


def merge_results(
    semantic_docs: List[Document],
    keyword_docs: List[Document],
    question: str,
    final_k: int = FINAL_K,
) -> List[Document]:
    by_id: Dict[str, Document] = {}
    combined_scores: Dict[str, float] = {}
    retrievals: Dict[str, set] = {}

    for doc in semantic_docs + keyword_docs:
        md = dict(doc.metadata or {})
        doc_id = str(md.get("id") or f"{md.get('source')}::{md.get('chunk')}")
        score = float(md.get("semantic_score", 0.0) or md.get("keyword_score", 0.0) or 0.0)

        if doc_id not in by_id:
            by_id[doc_id] = doc
            combined_scores[doc_id] = score
            retrievals[doc_id] = {md.get("retrieval", "")}
        else:
            combined_scores[doc_id] += score
            retrievals[doc_id].add(md.get("retrieval", ""))

    merged: List[Document] = []
    for doc_id, doc in by_id.items():
        md = dict(doc.metadata or {})
        md["retrievals"] = sorted(list(retrievals.get(doc_id, set())))
        md["combined_score"] = combined_scores.get(doc_id, 0.0)
        md["final_score"] = (
            float(md["combined_score"])
            + score_keyword_match(question, doc.page_content, md.get("source", ""))
            + (1.25 if len(md["retrievals"]) > 1 else 0.0)
        )
        merged.append(Document(page_content=doc.page_content, metadata=md))

    merged.sort(key=lambda d: float(d.metadata.get("final_score", 0.0)), reverse=True)

    final_docs: List[Document] = []
    seen = set()
    for doc in merged:
        key = (doc.metadata.get("source"), doc.metadata.get("chunk"))
        if key in seen:
            continue
        seen.add(key)
        final_docs.append(doc)
        if len(final_docs) >= final_k:
            break

    return final_docs


def expand_adjacent_docs(workspace_id: str, docs: List[Document]) -> List[Document]:
    index_data = get_index_data(workspace_id)
    chunk_lookup = index_data.get("chunk_lookup", {})
    if not docs or ADJACENT_CHUNKS <= 0:
        return docs

    expanded: List[Document] = []
    seen = set()

    for doc in docs:
        md = dict(doc.metadata or {})
        key = (md.get("source"), md.get("chunk"))
        if key not in seen:
            expanded.append(doc)
            seen.add(key)

        file_hash = str(md.get("file_hash"))
        chunk_idx = int(md.get("chunk", 0))

        for offset in range(-ADJACENT_CHUNKS, ADJACENT_CHUNKS + 1):
            if offset == 0:
                continue
            neighbor = chunk_lookup.get((file_hash, chunk_idx + offset))
            if not neighbor:
                continue
            nkey = (neighbor["source"], neighbor["chunk"])
            if nkey in seen:
                continue
            seen.add(nkey)
            expanded.append(
                Document(
                    page_content=neighbor["text"],
                    metadata={
                        "id": neighbor["id"],
                        "source": neighbor["source"],
                        "file_hash": neighbor["file_hash"],
                        "chunk": neighbor["chunk"],
                        "retrieval": "adjacent",
                        "retrievals": ["adjacent"],
                        "combined_score": 0.35,
                        "final_score": 0.35,
                    },
                )
            )

    return expanded[: max(FINAL_K * 2, 24)]


def build_context(docs: List[Document], max_chars: int) -> str:
    parts: List[str] = []
    used = 0

    for doc in docs:
        source = doc.metadata.get("source", "unknown")
        chunk = doc.metadata.get("chunk", 0)
        retrievals = doc.metadata.get("retrievals")
        retrieval = ",".join(retrievals) if retrievals else doc.metadata.get("retrieval", "")
        part = f"[source={source} chunk={chunk} retrieval={retrieval}] {doc.page_content}"
        part_len = len(part)

        if used + part_len > max_chars:
            remaining = max_chars - used
            if remaining > 250:
                parts.append(part[:remaining])
            break

        parts.append(part)
        used += part_len + 2

    return "\n\n".join(parts)


def rag_answer(workspace_id: str, question: str) -> str:
    query_variants = rewrite_queries(question)
    all_semantic: List[Document] = []
    all_keyword: List[Document] = []

    for q in query_variants:
        all_semantic.extend(semantic_search(workspace_id, q, k=SEMANTIC_K))
        all_keyword.extend(keyword_search(workspace_id, q, k=KEYWORD_K))

    docs = merge_results(all_semantic, all_keyword, question=question, final_k=FINAL_K)
    docs = expand_adjacent_docs(workspace_id, docs)
    context = build_context(docs, MAX_CONTEXT_CHARS)

    if not context.strip():
        return "I don't know based on the uploaded files."

    prompt = f"""{SYSTEM_PROMPT}

Context:
{context}

User question: {question}

Answer using only the context above. If the answer is partially supported, provide the strongest grounded answer possible and briefly say what remains unclear.
"""

    raw = generate_streaming(prompt)
    obj = safe_json_load(raw)
    answer = (obj.get("answer") or "").strip()
    if not answer:
        answer = "I don't know based on the uploaded files."

    if DEBUG_RETRIEVAL:
        st.session_state["last_debug_docs"] = [
            {
                "source": d.metadata.get("source"),
                "chunk": d.metadata.get("chunk"),
                "retrieval": d.metadata.get("retrieval"),
                "retrievals": d.metadata.get("retrievals", []),
                "preview": d.page_content[:400],
            }
            for d in docs
        ]

    touch_activity(workspace_id)
    return answer


# ============================
# MAIN APP
# ============================
def main() -> None:
    st.set_page_config(page_title=APP_TITLE, layout="centered")
    inject_styles()
    login_gate()

    workspace_id = ensure_tenant_context()

    if workspace_expired(workspace_id):
        clear_workspace_storage(workspace_id)
        st.session_state.messages = []
        st.session_state.indexed = False
        st.session_state.failed_files = []
        st.session_state.processed_files = []
        st.session_state.last_index_count = 0

    st.session_state.setdefault("messages", [])
    st.session_state.setdefault("indexed", False)
    st.session_state.setdefault("failed_files", [])
    st.session_state.setdefault("processed_files", [])
    st.session_state.setdefault("created_at", _now())
    st.session_state.setdefault("last_activity", _now())
    st.session_state.setdefault("last_debug_docs", [])
    st.session_state.setdefault("last_index_count", 0)

    index_data = get_index_data(workspace_id)
    if index_data.get("chunks"):
        st.session_state.indexed = True
        st.session_state.last_index_count = len(index_data.get("chunks", []))
        if not st.session_state.get("processed_files"):
            st.session_state.processed_files = [f.get("file_name", "") for f in index_data.get("files", [])]

    touch_activity(workspace_id)

    st.markdown(f'<div class="gothic-title">{APP_TITLE}</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="center-note">Ask Questions About Your Documents</div>',
        unsafe_allow_html=True,
    )

    if not st.session_state.indexed:
        st.markdown(
            '<div class="gothic-sub">Upload Documents and Ask Questions Instantly</div>',
            unsafe_allow_html=True,
        )

        uploaded = st.file_uploader(
            "Upload files",
            type=None,
            accept_multiple_files=True,
            label_visibility="visible",
        )

        st.markdown(
            f'<div class="small-muted">Up to {MAX_FILES} files • {MAX_FILE_MB}MB max per file • {MAX_TOTAL_UPLOAD_MB}MB total</div>',
            unsafe_allow_html=True,
        )

        c1, c2 = st.columns([2, 1])
        with c1:
            index_clicked = st.button("Upload & Prepare", type="primary", use_container_width=True)
        with c2:
            if st.button("Reset", use_container_width=True):
                reset_workspace()
                st.rerun()

        if index_clicked:
            ok, msg = validate_uploads(uploaded)
            if not ok:
                st.error(msg)
            else:
                try:
                    with st.spinner("Preparing your files..."):
                        count = index_files(workspace_id, uploaded)
                except Exception as e:
                    st.error(f"Indexing failed. Details: {e}")
                    return

                if count == 0:
                    st.warning(
                        "I couldn’t extract text from those files. Supported examples include PDF, DOC/DOCX, PPT/PPTX, ODT/ODS/ODP, TXT, CSV, XLS/XLSX/XLSB, Parquet, HTML, XML, YAML, JSON, JSONL, NDJSON, IPYNB, EML, MSG, RTF, RST, TEX, and code/text files."
                    )
                    if st.session_state.get("failed_files"):
                        st.error("Files that failed:\n\n" + "\n".join(f"- {x}" for x in st.session_state["failed_files"]))
                else:
                    st.success(
                        f"Ready. Indexed {count} chunks from {len(st.session_state.get('processed_files', []))} file(s)."
                    )
                    if st.session_state.get("failed_files"):
                        st.warning(
                            "Some files could not be processed:\n\n"
                            + "\n".join(f"- {x}" for x in st.session_state["failed_files"])
                        )
                    st.rerun()

        st.markdown(
            f"""
<div class="small-foot">
Files are automatically deleted after <strong>{IDLE_TTL_SECONDS // 60} minutes of inactivity</strong> or <strong>{WORKSPACE_TTL_SECONDS // 3600} hours maximum session time</strong>.<br/>
<em>This is a demo environment. Please upload sample documents only.</em><br/>
Questions or custom deployments: <strong>linkedin.com/in/thedannyscott</strong>
</div>
""",
            unsafe_allow_html=True,
        )
        return

    top = st.columns([3, 1])
    with top[0]:
        processed = st.session_state.get("processed_files", [])
        failed = st.session_state.get("failed_files", [])
        st.markdown(
            f'<div class="small-success">Indexed files: {len(processed)} • Failed files: {len(failed)} • Chunks: {st.session_state.get("last_index_count", 0)}</div>',
            unsafe_allow_html=True,
        )
    with top[1]:
        if st.button("Reset / New Upload", use_container_width=True):
            reset_workspace()
            st.rerun()

    if st.session_state.get("processed_files"):
        with st.expander("Show indexed file names"):
            st.write("\n".join(st.session_state["processed_files"]))

    if st.session_state.get("failed_files"):
        with st.expander("Show file processing issues"):
            st.write("\n".join(st.session_state["failed_files"]))

    if DEBUG_RETRIEVAL and st.session_state.get("last_debug_docs"):
        with st.expander("Debug: last retrieved chunks"):
            st.json(st.session_state["last_debug_docs"])

    for message in st.session_state.messages:
        avatar = ASSISTANT_AVATAR if message["role"] == "assistant" else None
        with st.chat_message(message["role"], avatar=avatar):
            st.markdown(message["content"])

    q = st.chat_input("Ask a question about your uploaded files…")
    if q:
        st.session_state.messages.append({"role": "user", "content": q})
        with st.chat_message("assistant", avatar=ASSISTANT_AVATAR):
            with st.spinner("Thinking..."):
                answer = rag_answer(workspace_id, q)
            st.markdown(answer)
        st.session_state.messages.append({"role": "assistant", "content": answer})


if __name__ == "__main__":
    main()
